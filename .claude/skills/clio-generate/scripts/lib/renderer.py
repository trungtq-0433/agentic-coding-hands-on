from __future__ import annotations

import base64
import copy
import os
import tempfile
from pathlib import Path
from typing import Dict
from typing import List
from typing import Optional
from typing import Tuple

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt

from lib.schemas import ContentItem
from lib.schemas import ContentType
from lib.schemas import SlideContent
from lib.templates import SVNProposalTemplate
from lib.templates.base import ContentKind, ShapeRole, SlideRoleConfig
from lib.profile_schema import ProjectProfile, section_has_content
from lib.shape_discovery import ShapeDiscovery



class PPTXRenderer:
    """Renderer for filling content into existing template slides"""

    # OOXML unit: 914400 EMU = 1 inch.  All other EMU constants derive from this.
    _EMU_PER_INCH: int = 914400

    # Rows are compressed to fit the safe area; this floor keeps text legible.
    _MIN_ROW_HEIGHT_EMU: int = int(0.25 * _EMU_PER_INCH)

    # `_estimate_max_chars` floors at 8 for a degenerately small box. At or
    # barely above that floor, truncating still cuts mid-word — a category
    # pill like "Intake & Registration" becomes the meaningless "Intake &…".
    # Below this threshold, skip truncation entirely and let auto-fit shrink
    # the full text instead (small-but-complete beats legible-but-meaningless).
    _MIN_MEANINGFUL_TRUNCATION_CHARS: int = 12

    # Minimum height for data rows to ensure content is readable.
    # At 8pt Noto Sans JP, one line ≈ 0.15"; two lines + padding ≈ 0.35".
    _MIN_DATA_ROW_HEIGHT_EMU: int = int(0.35 * _EMU_PER_INCH)

    # Font metrics for estimating table-cell wrap (8pt Noto Sans JP). Shared by
    # the pre-fill split-size estimate and the post-fill row-height application
    # so both agree on what a row actually needs.
    _TABLE_FONT_PT: int = 8
    _TABLE_CJK_W_INCH: float = _TABLE_FONT_PT / 72          # ~0.111" per CJK char
    _TABLE_LATIN_W_INCH: float = _TABLE_FONT_PT * 0.6 / 72  # ~0.067" per Latin char
    _TABLE_LINE_H_EMU: int = int(_TABLE_FONT_PT * 1.6 / 72 * _EMU_PER_INCH)
    _TABLE_CELL_PADDING_EMU: int = int(0.07 * _EMU_PER_INCH)

    def __init__(self):
        self.template_config = None
        self._output_dir: Optional[Path] = None
        # Safe content area bottom — computed dynamically from slide height in render_from_profile.
        # Fallback: 5" (suits a 5.625" widescreen slide with footer at ~5.25").
        self._SVN_SAFE_BOTTOM_EMU: int = int(5.0 * self._EMU_PER_INCH)
        # Actual slide width — set from prs.slide_width in render_from_profile.
        self._slide_width_emu: int = 0
        # Text-role items too long for their fixed template box: truncated in
        # place, full text queued here as {(section, role_name): [full_text, ...]}
        # and rendered as an appendix extra-slide after the main fill pass —
        # reset per render_from_profile call.
        self._overflow_extras: Dict[Tuple[str, str], List[str]] = {}

    def render(
        self,
        slides: List[SlideContent],
        template: str,
        output_name: str,
        output_dir: Optional[str] = None,
    ) -> Path:
        """
        Fill content into existing template slides

        Args:
            slides: List of slide content to fill
            template: Template PPTX file path (e.g., 'SVN Proposal Menu.pptx')
            output_name: Output file name (without .pptx extension)
            output_dir: Deprecated/ignored. Output directory is resolved from env `SLIDE_GENERATOR__OUTPUTS_PATH`.

        Returns:
            Path to generated PPTX file
        """
        print(f'Filling content into template slides: {output_name}')

        # Resolve output_dir once so image path resolution can use it
        self._output_dir = self._resolve_output_dir(output_dir)

        # Auto-detect template config from template name
        self._detect_template_config(template)

        # Load template presentation
        prs = self._load_template(template)

        # Fill content into existing slides
        self._fill_existing_slides(prs, slides)

        # Save presentation
        output_file = self._save_presentation(prs, output_name, output_dir)

        return output_file

    @staticmethod
    def _resolve_output_dir(output_dir: Optional[str]) -> Path:
        """Resolve output directory path.

        Priority: explicit output_dir > env SLIDE_GENERATOR__OUTPUTS_PATH > CWD/outputs
        """
        if output_dir:
            return Path(output_dir)
        env_outputs = os.getenv('SLIDE_GENERATOR__OUTPUTS_PATH')
        if env_outputs:
            return Path(env_outputs).expanduser()
        return Path.cwd() / 'outputs'

    def _detect_template_config(self, template: str):
        """Auto-detect and load template config from template name"""
        if 'svn' in template.lower():
            self.template_config = SVNProposalTemplate()
            print(f'Auto-detected SVN template config: {len(self.template_config.slide_configs)} slides configured')
        else:
            self.template_config = None
            print('No specific template config detected, using default rendering')

    def _load_template(self, template: str) -> Presentation:
        """Load template presentation — supports absolute paths and relative names."""
        template_path = Path(template)
        if template_path.is_absolute() and template_path.exists():
            prs = Presentation(str(template_path))
            print(f'Loaded template: {template_path.name} ({len(prs.slides)} slides)')
            return prs
        rel_path = Path(__file__).parent / 'templates' / template
        if rel_path.exists():
            prs = Presentation(str(rel_path))
            print(f'Loaded template: {template} from {rel_path} ({len(prs.slides)} slides)')
            return prs
        print(f'Template not found: {template} at {template_path} or {rel_path}')
        return Presentation()

    def _save_presentation(self, prs: Presentation, output_name: str, output_dir: Optional[str] = None) -> Path:
        """Save presentation to output directory.

        Priority: explicit output_dir > env SLIDE_GENERATOR__OUTPUTS_PATH > CWD/outputs
        """
        output_path = self._output_dir or self._resolve_output_dir(output_dir)

        try:
            output_path.mkdir(parents=True, exist_ok=True)
        except PermissionError as e:
            # Make error clearer for API layer
            raise PermissionError(f'Permission denied creating outputs directory: {output_path}') from e
        output_file = output_path / f'{output_name}.pptx'

        prs.save(str(output_file))
        print(f'Saved PPTX to: {output_file}')

        return output_file

    def _fill_existing_slides(self, prs: Presentation, slides_content: List[SlideContent]):
        """Fill content into existing slides"""
        print(f'Filling content into {len(slides_content)} existing slides')

        # Sort by slide number so that the offset tracking works correctly when
        # overflow slides are inserted in the middle of the presentation.
        sorted_contents = sorted(slides_content, key=lambda s: s.slide_number)

        # Tracks how many extra slides have been inserted before the current one.
        slide_offset = 0

        for slide_content in sorted_contents:
            slide_num = slide_content.slide_number

            # Check if slide is protected
            if self.template_config and not self.template_config.should_fill_slide(slide_num):
                print(f'Slide {slide_num} is protected - skipping')
                continue

            # Adjusted index accounts for any extra slides inserted before this slide.
            adjusted_index = slide_num - 1 + slide_offset

            # Check if slide exists
            if adjusted_index < len(prs.slides):
                slide = prs.slides[adjusted_index]
                print(f'Filling content into slide {slide_num} (template index: {adjusted_index})')

                # Fill content; if a table overflows it will insert extra slides and
                # return the count so we can update the offset.
                extra = self._fill_with_overflow_handling(prs, slide, slide_content, adjusted_index)
                slide_offset += extra
            else:
                print(
                    f'Slide {slide_num} not found in template '
                    f'(total slides: {len(prs.slides)}, adjusted index: {adjusted_index})',
                )

    # ------------------------------------------------------------------
    # Table-overflow helpers
    # ------------------------------------------------------------------

    def _fill_with_overflow_handling(
        self,
        prs: Presentation,
        slide,
        slide_content: SlideContent,
        slide_index: int,
    ) -> int:
        """Fill a slide, splitting into multiple slides when a table is too tall
        to fit within the safe content area above the footer band.

        Returns the number of *extra* slides inserted (0 when no overflow).
        """
        shape_contents = getattr(slide_content, 'shape_contents', {}) or {}

        # Find the first table-type shape content for this slide
        table_key, table_data = self._find_table_content(slide_content, shape_contents)

        if not table_key:
            # No table present – fill normally
            self._fill_slide_content(slide, slide_content)
            return 0

        # Legacy path: shape position is unknown here, so estimate from safe area.
        # Assume content starts ~1.5" below the top (title band); minimum row
        # height gives the theoretical maximum rows that could ever fit.
        data_area_est = self._SVN_SAFE_BOTTOM_EMU - int(1.5 * self._EMU_PER_INCH)
        legacy_max_rows = max(1, data_area_est // self._MIN_ROW_HEIGHT_EMU)
        chunks = self._split_table_data(table_data, legacy_max_rows)

        if len(chunks) <= 1:
            # Fits on one slide – fill normally
            self._fill_slide_content(slide, slide_content)
            return 0

        print(
            f'  - Slide {slide_content.slide_number}: table split '
            f'({len(chunks)} page(s), up to {legacy_max_rows} rows/page)',
        )

        # Fill the original slide with the first chunk
        first_content = self._build_chunked_slide_content(
            slide_content, table_key, chunks[0], is_first=True,
        )
        self._fill_slide_content(slide, first_content)

        # Create one continuation slide per remaining chunk
        extra_slides = 0
        current_index = slide_index
        for i, chunk in enumerate(chunks[1:], 1):
            new_slide = self._duplicate_slide(prs, current_index)
            extra_slides += 1
            current_index += 1

            cont_content = self._build_chunked_slide_content(
                slide_content, table_key, chunk, is_first=False,
            )
            self._fill_slide_content(new_slide, cont_content)
            print(
                f'  - Added continuation slide {i}/{len(chunks) - 1} '
                f'at template position {current_index + 1}',
            )

        return extra_slides

    def _find_table_content(
        self,
        slide_content: SlideContent,
        shape_contents: dict,
    ):
        """Return (content_key, markdown_table_string) for the first table-type
        shape found in this slide's config.  Returns (None, None) if none found.
        """
        config = (
            self.template_config.get_slide_config(slide_content.slide_number)
            if self.template_config
            else None
        )
        if not config or not config.shape_targets:
            return None, None

        for target in config.shape_targets:
            content = shape_contents.get(target.content_key, '')
            if content and self._is_table_data(content):
                return target.content_key, content

        return None, None

    def _split_table_data(self, table_data: str, max_rows: int) -> List[str]:
        """Split a markdown table string into chunks of at most *max_rows* data rows.

        Each chunk keeps the original header row(s) and the separator row so that
        each chunk is a valid, self-contained markdown table.

        Returns a list of markdown table strings (length 1 when no split needed).
        """
        lines = [line for line in table_data.strip().split('\n') if line.strip()]
        if not lines:
            return [table_data]

        # Separate header / separator lines from data lines
        header_lines: List[str] = []
        data_lines: List[str] = []
        separator_found = False

        for line in lines:
            stripped = line.strip()
            is_separator = (
                not separator_found
                and '|' in stripped
                and all(c in '|-: \t' for c in stripped)
            )
            if is_separator:
                header_lines.append(line)
                separator_found = True
            elif not separator_found:
                header_lines.append(line)
            else:
                data_lines.append(line)

        if len(data_lines) <= max_rows:
            return [table_data]

        header_block = '\n'.join(header_lines)
        chunks: List[str] = []
        for i in range(0, len(data_lines), max_rows):
            row_block = '\n'.join(data_lines[i : i + max_rows])
            chunks.append(f'{header_block}\n{row_block}')

        return chunks

    def _build_chunked_slide_content(
        self,
        original: SlideContent,
        table_key: str,
        table_chunk: str,
        is_first: bool,
    ) -> SlideContent:
        """Return a new SlideContent with *table_key* replaced by *table_chunk*.

        For continuation slides (``is_first=False``) all non-table text shapes
        are cleared so description text is not repeated across slides.
        """
        new_shape_contents = dict(getattr(original, 'shape_contents', {}) or {})
        new_shape_contents[table_key] = table_chunk

        if not is_first:
            config = (
                self.template_config.get_slide_config(original.slide_number)
                if self.template_config
                else None
            )
            if config and config.shape_targets:
                for target in config.shape_targets:
                    key = target.content_key
                    if key != table_key and key in new_shape_contents:
                        if not self._is_table_data(new_shape_contents[key]):
                            new_shape_contents[key] = ''  # Clear description text

        return SlideContent(
            slide_number=original.slide_number,
            layout=original.layout,
            title=original.title,
            subtitle=original.subtitle,
            content=original.content,
            shape_contents=new_shape_contents,
        )

    @staticmethod
    def _ensure_unique_slide_partname(prs: Presentation, slide) -> None:
        """Work around a python-pptx limitation: ``Slides.add_slide()``
        assigns the new slide's partname from ``len(sldIdLst) + 1`` — a
        naive count of the CURRENT slide list — not from the set of
        partnames actually in use. Once slides have been deleted out of
        order (as ``_prune_and_reorder`` does via manual XML surgery,
        working around python-pptx having no public slide-delete API),
        that count drifts below the real max partname index and collides
        with an original template slide that's still alive. python-pptx's
        zip writer doesn't error on this — it silently emits two different
        slide XML parts under the same path, and whichever is written last
        wins for BOTH relationships that point at that name (the earlier
        one's content becomes unreachable). Call right after every
        ``add_slide()`` to reassign a name checked against every part
        actually in the package.
        """
        from pptx.opc.package import PackURI

        package = prs.part.package
        tmpl = '/ppt/slides/slide%d.xml'
        prefix = tmpl[: (tmpl % 42).find('42')]
        others = {p.partname for p in package.iter_parts()
                  if p.partname.startswith(prefix) and p is not slide.part}
        if slide.part.partname in others:
            n = len(others) + 1
            while tmpl % n in others:
                n += 1
            slide.part.partname = PackURI(tmpl % n)

    def _duplicate_slide(self, prs: Presentation, slide_index: int):
        """Duplicate the slide at *slide_index* and insert the copy immediately
        after it in the presentation.

        Returns the new slide object (at position ``slide_index + 1``).
        """
        source_slide = prs.slides[slide_index]

        # 1. Add a new blank slide with the same layout – this is appended at the end.
        new_slide = prs.slides.add_slide(source_slide.slide_layout)
        self._ensure_unique_slide_partname(prs, new_slide)

        # 2. Replace the new slide's shape tree with a deep copy of the source's.
        source_spTree = source_slide.shapes._spTree
        new_spTree = new_slide.shapes._spTree

        for child in list(new_spTree):
            new_spTree.remove(child)
        for child in source_spTree:
            new_child = copy.deepcopy(child)
            self._remap_image_relationships(source_slide.part, new_slide.part, new_child)
            new_spTree.append(new_child)

        # 3. Move the new slide from the end to slide_index + 1 by reordering
        #    the <p:sldIdLst> in the presentation XML.
        pml_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        try:
            presentation_elem = prs.part._element
            sldIdLst = presentation_elem.find(f'{{{pml_ns}}}sldIdLst')
            if sldIdLst is not None:
                sld_id_elems = list(sldIdLst)
                if sld_id_elems:
                    new_sldId = sld_id_elems[-1]          # just appended → at end
                    source_sldId = sld_id_elems[slide_index]  # source position
                    sldIdLst.remove(new_sldId)
                    source_sldId.addnext(new_sldId)
                    print(
                        f'  - Duplicated slide at index {slide_index}, '
                        f'inserted at index {slide_index + 1}',
                    )
        except Exception as e:
            print(f'  - Error repositioning duplicated slide: {e}')

        return prs.slides[slide_index + 1]

    @staticmethod
    def _remove_shape(shape) -> None:
        """Remove a shape entirely (not just its text) from its slide's shape tree."""
        if shape is None:
            return
        elem = shape._element
        parent = elem.getparent()
        if parent is not None:
            parent.remove(elem)

    @staticmethod
    def _remove_decorations_around(slide, ref_shapes: list, padding_emu: int = 228600) -> int:
        """Remove sibling shapes (card background, icon badge, icon picture)
        that visually belong to an unused title/body (or text-item) slot but
        aren't tracked by the role config, so blanking the tracked shapes'
        text alone would leave an empty-looking card/box behind.

        Purely geometric — the role config only knows about title_shape/
        body_shape, not per-slot decoration, and this template's decoration
        shapes aren't at any predictable index/name offset from their title
        (e.g. an icon picture can sit anywhere else in the shape list).
        Any top-level shape whose bounding box is contained within the union
        of `ref_shapes`'s boxes, expanded by `padding_emu` (default 0.25" —
        covers a card background that bleeds slightly past its title+body,
        without reaching a neighboring card across this template's ~0.3"+
        gaps), is removed. `ref_shapes` themselves are excluded.
        """
        boxes = [
            (s.left, s.top, s.left + s.width, s.top + s.height)
            for s in ref_shapes if s is not None and s.left is not None
        ]
        if not boxes:
            return 0
        # Compare by the underlying XML element, not the shape wrapper's own
        # id() — python-pptx returns a fresh wrapper instance each time
        # `.shapes` is iterated/indexed, so two wrappers for the SAME shape
        # never share an id() even though `sh1._element is sh2._element`.
        # Using id(wrapper) here would make every ref_shape match itself as
        # "contained within its own bounding box" and get deleted too.
        exclude_ids = {id(s._element) for s in ref_shapes if s is not None}
        l = min(b[0] for b in boxes) - padding_emu
        t = min(b[1] for b in boxes) - padding_emu
        r = max(b[2] for b in boxes) + padding_emu
        b = max(b[3] for b in boxes) + padding_emu
        removed = 0
        for sh in list(slide.shapes):
            if id(sh._element) in exclude_ids or sh.left is None:
                continue
            sl, st = sh.left, sh.top
            sr, sb = sl + sh.width, st + sh.height
            if sl >= l and st >= t and sr <= r and sb <= b:
                PPTXRenderer._remove_shape(sh)
                removed += 1
        return removed

    @staticmethod
    def _remap_image_relationships(source_part, new_part, element) -> None:
        """Fix up `r:embed` references after deep-copying XML between parts.

        A deep-copied <p:pic>/<a:blip> keeps its source slide's relationship
        ID (e.g. "rId1") verbatim, but that ID is local to the source slide's
        own .rels — the destination slide (freshly created via `add_slide`,
        or the same slide the element is merely being duplicated within)
        doesn't necessarily have a matching relationship at that ID. Left
        unfixed, the reference silently dangles and the picture renders as
        nothing (no error, just a blank image) in PowerPoint/LibreOffice.
        """
        if source_part is new_part:
            return
        from pptx.oxml.ns import qn
        r_embed = qn('r:embed')
        for el in element.iter():
            old_rid = el.get(r_embed)
            if not old_rid:
                continue
            try:
                image_part = source_part.related_part(old_rid)
                reltype = source_part.rels[old_rid].reltype
                new_rid = new_part.relate_to(image_part, reltype)
                el.set(r_embed, new_rid)
            except Exception:
                pass

    @classmethod
    def _clone_shape(cls, slide, source_shape, top: Optional[int] = None):
        """Deep-copy *source_shape* onto *slide* (same slide or another), optionally
        repositioning the copy's top. Returns the new Shape."""
        new_elem = copy.deepcopy(source_shape._element)
        source_part = source_shape.part
        cls._remap_image_relationships(source_part, slide.part, new_elem)
        slide.shapes._spTree.append(new_elem)
        new_shape = slide.shapes[-1]
        if top is not None:
            new_shape.top = top
        return new_shape

    # ------------------------------------------------------------------

    def _fill_slide_content(self, slide, slide_content: SlideContent):
        """Fill content into an existing slide"""
        slide_num = slide_content.slide_number
        config = self.template_config.get_slide_config(slide_num) if self.template_config else None

        # Strategy 1: Shape targeting (if shape_contents exists)
        if self._has_shape_contents(slide_content) and config and config.shape_targets:
            print(f'  - Using shape targeting for slide {slide_num}')
            self._fill_with_shape_targets(slide, slide_content, config)
            return

        # Strategy 2: Default filling
        self._fill_slide_default(slide, slide_content)

    def _has_shape_contents(self, slide_content: SlideContent) -> bool:
        """Check if slide_content has shape_contents attribute"""
        return bool(
            hasattr(slide_content, '__dict__') and
            'shape_contents' in slide_content.__dict__ and
            slide_content.shape_contents,
        )

    def _fill_with_shape_targets(self, slide, slide_content, config):
        """Fill content using shape targets from config"""
        shape_contents = getattr(slide_content, 'shape_contents', {})

        if not shape_contents:
            print('  - No shape_contents found in slide_content')
            return

        print(f'  - Found {len(shape_contents)} shape contents to fill')

        # Update title if exists
        self._update_slide_title(slide, slide_content.title)

        # Fill each shape target
        for target in config.shape_targets:
            self._fill_single_shape_target(slide, target, shape_contents, config)

    def _update_slide_title(self, slide, title: Optional[str]):
        """Update slide title if exists"""
        if slide.shapes.title and title:
            slide.shapes.title.text = title
            print(f'  - Updated title: {title}')

    def _fill_single_shape_target(self, slide, target, shape_contents: Dict[str, str], config=None):
        """Fill a single shape target with content"""
        content_key = target.content_key

        # Check if content exists for this key
        if content_key not in shape_contents:
            print(f"  - Content key '{content_key}' not found in markdown")
            return

        content = shape_contents[content_key]

        # Find shape by index or name
        shape = self._find_shape(slide, target)

        if not shape:
            print(f'  - Shape not found for target: {content_key}')
            return

        # Detect content type and fill accordingly
        if self._is_image_path(content):
            # Content is an image path
            self._replace_shape_with_image(slide, shape, content, target)
        elif target.fill_cols or self._is_table_data(content):
            # fill_cols always implies a table target (even single-row);
            # otherwise fall back to markdown table detection.
            self._fill_table(shape, content, target)
        elif shape.has_text_frame:
            # Content is text
            shape.text_frame.text = content
            self._apply_text_font(shape.text_frame)
            
            # Apply 10pt line spacing for slide 4
            if config and config.slide_number == 4:
                self._apply_line_spacing(shape.text_frame, Pt(10))
                print(f"  - Applied 10pt line spacing to slide 4 shape '{content_key}'")
            
            print(f"  Filled '{content_key}' into {shape.name} ({len(content)} chars)")
        else:
            print(f'  - Shape {shape.name} has no text frame and content is not an image or table')

    def _is_image_path(self, content: str) -> bool:
        """Check if content is an image file path or base64 data URI"""
        content = content.strip()
        # Check base64 data URI: data:image/...;base64,...
        if content.startswith('data:image/'):
            return True
        # Check if it's a file path with image extension
        image_extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.svg']
        return any(content.lower().endswith(ext) for ext in image_extensions)

    def _decode_base64_image(self, data_uri: str) -> Optional[str]:
        """Decode a base64 data URI to a temporary file and return the file path.

        Supports format: data:image/<ext>;base64,<data>

        Returns:
            Absolute path to the temporary file, or None on error.
        """
        try:
            # Parse header: data:image/png;base64,<data>
            header, encoded = data_uri.split(',', 1)
            # Extract mime type, e.g. image/png
            mime_part = header.split(':')[1].split(';')[0]  # e.g. "image/png"
            ext = '.' + mime_part.split('/')[1]  # e.g. ".png"
            # Handle special cases
            if ext == '.jpeg':
                ext = '.jpg'
            elif ext == '.svg+xml':
                ext = '.svg'

            image_bytes = base64.b64decode(encoded)
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
            tmp.write(image_bytes)
            tmp.flush()
            tmp.close()
            print(f'  - Decoded base64 image to temp file: {tmp.name} ({len(image_bytes)} bytes)')
            return tmp.name
        except Exception as e:
            print(f'  - Failed to decode base64 image: {e}')
            return None

    def _replace_shape_with_image(self, slide, old_shape, image_path: str, target):
        """Replace a shape with an image (supports file path or base64 data URI)"""
        tmp_file: Optional[str] = None
        try:
            # Get shape position and size
            left = old_shape.left
            top = old_shape.top
            width = old_shape.width
            height = old_shape.height

            image_path = image_path.strip()

            # Handle base64 data URI
            if image_path.startswith('data:image/'):
                tmp_file = self._decode_base64_image(image_path)
                if not tmp_file:
                    print(f"  - Cannot decode base64 image for shape '{target.content_key}'")
                    return
                image_path = tmp_file
            else:
                # Resolve file path
                if not Path(image_path).is_absolute():
                    rel_path = Path(image_path)
                    if rel_path.exists():
                        image_path = str(rel_path)
                    elif self._output_dir:
                        resolved = self._output_dir / image_path
                        if resolved.exists():
                            image_path = str(resolved)
                    else:
                        outputs_path = Path('outputs') / image_path
                        if outputs_path.exists():
                            image_path = str(outputs_path)

                # Check if image exists
                if not Path(image_path).exists():
                    print(f'  - Image not found: {image_path}')
                    return

            # Add new image at same position (overlays the placeholder shape)
            try:
                slide.shapes.add_picture(
                    str(image_path),
                    left, top,
                    width=width,
                    height=height,
                )
                print(f"  Replaced shape '{target.content_key}' with image: {Path(image_path).name}")
            except Exception as e:
                print(f'  - Error adding image: {e}')

        except Exception as e:
            print(f'  - Error replacing shape with image: {e}')
        finally:
            # Clean up temp file if created from base64
            if tmp_file and Path(tmp_file).exists():
                try:
                    Path(tmp_file).unlink()
                    print(f'  - Cleaned up temp file: {tmp_file}')
                except Exception:
                    pass

    def _is_table_data(self, content: str) -> bool:
        """Check if content is table data (markdown table format).

        Accepts two formats:
        1. Standard markdown table with header + separator:
               | col1 | col2 |
               |------|------|
               | val1 | val2 |
        2. No-header row-only table (used for fill_cols targeting):
               | val1 |
               | val2 |
               | val3 |
        """
        content = content.strip()
        lines = [line.strip() for line in content.split('\n') if line.strip()]

        if len(lines) < 2:
            return False

        # All lines must contain at least one pipe
        if not all('|' in line for line in lines):
            return False

        # Format 1: standard table with separator line
        for line in lines[1:3]:
            if all(c in '|-: \t' for c in line):
                return True

        # Format 2: no-separator — every line starts and ends with | (row-only table)
        if all(line.startswith('|') and line.endswith('|') for line in lines):
            return True

        return False

    def _fill_table(self, shape, table_data: str, target):
        """Fill an existing table shape with data from markdown table format"""
        try:
            # Check if shape has table
            if not shape.has_table:
                print(f"  - Shape '{target.content_key}' is not a table")
                return

            table = shape.table
            rows_data = self._parse_md_table_rows(table_data)

            if not rows_data:
                print('  - No valid table data found')
                return

            # Get current table dimensions
            current_rows = len(table.rows)
            current_cols = len(table.columns)
            needed_rows = len(rows_data)
            needed_cols = len(rows_data[0]) if rows_data else 0

            # Adjust table size if needed
            if needed_rows > current_rows:
                # Add rows if we need more
                rows_to_add = needed_rows - current_rows
                print(f'  - Adding {rows_to_add} rows to table (current: {current_rows}, needed: {needed_rows})')
                self._add_table_rows(table, rows_to_add)
                # After adding rows, update current_rows to reflect the new count.
                # Note: python-pptx's table.rows collection is live, so this should
                # reflect the new size immediately, but we recalculate for safety.
                current_rows = len(table.rows)
            elif needed_rows < current_rows:
                # Fewer rows needed - remove the extra rows from the XML so they
                # don't appear as empty rows in the rendered slide.
                extra_rows = current_rows - needed_rows
                print(
                    '  - Table has %s rows but only %s rows of data '
                    '(removing %s unused rows)',
                    current_rows,
                    needed_rows,
                    extra_rows,
                )
                self._remove_table_rows(table, needed_rows)
                current_rows = len(table.rows)

            if needed_cols > current_cols:
                print(f'  - Data has {needed_cols} columns but table only has {current_cols} columns')
                needed_cols = current_cols  # Limit to available columns
            elif needed_cols < current_cols and not target.fill_cols:
                print(f'  - Data has {needed_cols} columns but table has {current_cols} columns '
                      f'(clearing {current_cols - needed_cols} unmatched column(s) instead of '
                      f'leaving stale template placeholder text)')

            # Fill table cells
            filled_count = 0
            for row_idx, row_data in enumerate(rows_data):
                # After adding rows via XML, we can fill up to needed_rows
                # Use try-except to handle any edge cases where table structure might not match
                try:
                    # If fill_cols is specified, map row_data positionally onto those columns;
                    # columns not listed are left untouched (preserving template content).
                    if target.fill_cols:
                        col_pairs = list(zip(target.fill_cols, row_data))
                    else:
                        col_pairs = list(enumerate(row_data))

                    filled_col_idxs = set()
                    for col_idx, cell_value in col_pairs:
                        if col_idx >= current_cols:
                            break  # Skip extra columns in data

                        cell = table.cell(row_idx, col_idx)
                        cell.text = cell_value
                        self._apply_table_cell_font(cell)
                        filled_count += 1
                        filled_col_idxs.add(col_idx)

                    # Without fill_cols, the data is expected to cover every
                    # column; any column the data doesn't reach gets cleared
                    # so stale template placeholder text can't leak through.
                    # fill_cols mode is exempt — it deliberately targets a
                    # subset of columns and leaves the rest (e.g. a fixed
                    # label column) untouched by design.
                    if not target.fill_cols:
                        for col_idx in range(current_cols):
                            if col_idx not in filled_col_idxs:
                                table.cell(row_idx, col_idx).text = ''
                except (IndexError, AttributeError) as e:
                    print(f'  - Cannot fill row {row_idx}: {e}')
                    break  # Stop filling if we hit an error

            print(f"  Filled '{target.content_key}' table: {len(rows_data)} rows x {needed_cols} cols ({filled_count} cells)")
            self._normalize_table_col_widths(table)
            self._apply_content_row_heights(table)
            self._fit_table_to_bounds(shape)

        except Exception as e:
            print(f'  - Error filling table: {e}')

    def _add_table_rows(self, table, num_rows: int):
        """Add rows to an existing table by manipulating XML

        python-pptx uses lxml internally, so we can manipulate the XML directly
        to add new rows by cloning the structure of the last row.
        """
        try:
            from lxml import etree

            # Get the table XML element (this is an lxml element)
            tbl = table._tbl

            # Get the number of existing rows
            current_row_count = len(table.rows)

            # Get the last row to use as template for new rows
            if current_row_count == 0:
                print('  - Cannot add rows: table has no existing rows to use as template')
                return

            # Use positive index instead of negative index (table.rows doesn't support negative indexing)
            last_row = table.rows[current_row_count - 1]
            last_row_xml = last_row._tr

            # Define namespace for XPath queries
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

            # Add new rows by cloning the last row
            for i in range(num_rows):
                # Create a deep copy of the last row by serializing and parsing
                new_row_xml = etree.fromstring(etree.tostring(last_row_xml))

                # Clear text content in all cells of the new row
                for tc in new_row_xml.xpath('.//a:tc', namespaces=ns):
                    txBody = tc.find('.//a:txBody', namespaces=ns)
                    if txBody is not None:
                        # Get all paragraphs
                        paragraphs = txBody.findall('.//a:p', namespaces=ns)
                        if paragraphs:
                            # Clear text in first paragraph
                            first_p = paragraphs[0]
                            for t_elem in first_p.xpath('.//a:t', namespaces=ns):
                                t_elem.text = ''
                            # Remove other paragraphs (keep structure but remove extra content)
                            for p in paragraphs[1:]:
                                txBody.remove(p)

                # Append the new row to the table
                tbl.append(new_row_xml)

            print(f'  - Successfully added {num_rows} rows to table')

        except ImportError:
            print('  - lxml not available. python-pptx requires lxml, please install it: pip install lxml')
        except Exception as e:
            print(f'  - Error adding rows to table: {e}')
            import traceback
            print(f'  - Traceback: {traceback.format_exc()}')

    def _remove_table_rows(self, table, keep_rows: int):
        """Remove all rows after index *keep_rows* from the table by manipulating XML.

        Args:
            table: python-pptx Table object
            keep_rows: Number of rows to keep (rows at index 0 … keep_rows-1 are retained).
        """
        try:
            from lxml import etree  # noqa: F401 (just to confirm availability)

            tbl = table._tbl
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            tr_elements = tbl.findall('a:tr', namespaces=ns)

            rows_to_remove = tr_elements[keep_rows:]
            for tr in rows_to_remove:
                tbl.remove(tr)

            print(
                f'  - Removed {len(rows_to_remove)} extra row(s); '
                f'table now has {keep_rows} row(s)',
            )
        except Exception as e:
            print(f'  - Error removing table rows: {e}')

    @staticmethod
    def _is_cjk_char(ch: str) -> bool:
        return '　' <= ch <= '鿿' or '＀' <= ch <= '￯'

    def _estimate_row_content_height_emu(
        self, cells: List[str], col_widths_inch: List[float],
    ) -> int:
        """Estimate the EMU height *cells* need to render without further wrap,
        given each column's width. Shared by the pre-fill split-size estimate
        (_max_rows_for_shape) and the post-fill row-height pass
        (_apply_content_row_heights) so both agree on what a row actually needs.
        """
        max_lines = 1
        for j, text in enumerate(cells):
            if j >= len(col_widths_inch) or not text:
                continue
            col_w = col_widths_inch[j]
            text_width = sum(
                self._TABLE_CJK_W_INCH if self._is_cjk_char(ch) else self._TABLE_LATIN_W_INCH
                for ch in text
            )
            lines = max(1, int(text_width / col_w) + 1)
            max_lines = max(max_lines, lines)
        return int(max_lines * self._TABLE_LINE_H_EMU + self._TABLE_CELL_PADDING_EMU)

    @staticmethod
    def _parse_md_table_rows(table_data: str) -> List[List[str]]:
        """Parse a markdown pipe-table string into row cell-lists (header first).

        Drops separator lines (e.g. ``|---|---|``). Shared by `_fill_table` (which
        writes the cells) and `_max_rows_for_shape` (which estimates how many rows
        of this exact data will fit) so both operate on the same parsed rows.
        """
        lines = [line.strip() for line in table_data.strip().split('\n') if line.strip()]
        rows_data = []
        for line in lines:
            if all(c in '|-: \t' for c in line):
                continue  # separator line
            cells = [cell.strip() for cell in line.split('|')]
            cells = [cell for cell in cells if cell]
            if cells:
                rows_data.append(cells)
        return rows_data

    def _apply_content_row_heights(self, table) -> None:
        """Bump data-row heights so each row is tall enough for its actual text content.

        The template may have very compact rows (e.g. 0.262") designed for short
        placeholder text.  When real content is longer and word-wraps, some renderers
        (LibreOffice, Google Slides preview) auto-expand rows — causing the table to
        grow past the footer even though the python-pptx shape is technically in-bounds.

        This method estimates the required row height from text length + column width,
        then sets each data row to max(template_height, content_height, MIN_DATA_ROW).
        After this call, _fit_table_to_bounds sees the true height and can compress
        or trigger a split correctly.
        """
        try:
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            emu = self._EMU_PER_INCH

            cols = list(table.columns)
            col_widths_inch = [max(c.width / emu, 0.5) for c in cols] if cols else []

            tbl = table._tbl
            tr_elements = tbl.findall('a:tr', namespaces=ns)
            changed = 0

            for idx, tr in enumerate(tr_elements):
                if idx == 0:
                    continue  # skip header row

                current_h = int(tr.get('h', '0'))
                min_h = max(current_h, self._MIN_DATA_ROW_HEIGHT_EMU)

                cells = []
                for j, tc in enumerate(tr.findall('a:tc', namespaces=ns)):
                    if j >= len(col_widths_inch):
                        break
                    texts = tc.findall('.//a:t', namespaces=ns)
                    cells.append(''.join(t.text or '' for t in texts).strip())

                content_h = self._estimate_row_content_height_emu(cells, col_widths_inch)
                new_h = max(min_h, content_h)

                if new_h != current_h:
                    tr.set('h', str(new_h))
                    changed += 1

            if changed:
                print(f'  - Content-adjusted {changed} row height(s) '
                      f'(min {self._MIN_DATA_ROW_HEIGHT_EMU/emu:.3f}")')
        except Exception as e:
            print(f'  - _apply_content_row_heights error: {e}')

    def _max_rows_for_shape(
        self, shape, rows_data: Optional[List[List[str]]] = None,
    ) -> int:
        """Compute max data rows that fit above the footer for *shape*.

        When *rows_data* (the actual cell text about to be filled, header
        excluded) is given, each row's real wrap-estimated height is used — the
        same heuristic `_apply_content_row_heights` applies after filling — so
        the split point matches what will actually be rendered instead of
        assuming every row is as short as the template's placeholder text.
        Without *rows_data*, falls back to the template row heights (each
        bumped to at least `_MIN_DATA_ROW_HEIGHT_EMU`) plus the last template
        row's height as the per-new-row cost.
        """
        emu = self._EMU_PER_INCH
        default_row_h = emu // 3  # ~0.33" when a row has no explicit height

        try:
            avail_height = max(0, self._SVN_SAFE_BOTTOM_EMU - shape.top)

            if not shape.has_table:
                return max(1, avail_height // self._MIN_ROW_HEIGHT_EMU)

            table = shape.table
            all_rows = list(table.rows)
            if not all_rows:
                return max(1, avail_height // self._MIN_ROW_HEIGHT_EMU)

            header_h = all_rows[0].height or default_row_h
            avail_for_data = max(0, avail_height - header_h)

            if rows_data:
                col_widths_inch = [max(c.width / emu, 0.5) for c in table.columns]
                result = 0
                used = 0
                for cells in rows_data:
                    row_h = max(
                        self._MIN_DATA_ROW_HEIGHT_EMU,
                        self._estimate_row_content_height_emu(cells, col_widths_inch),
                    )
                    if result and used + row_h > avail_for_data:
                        break
                    used += row_h
                    result += 1

                print(f'  - max_rows_for_shape (content-aware): avail={avail_height/emu:.2f}" '
                      f'header={header_h/emu:.2f}" data={used/emu:.2f}" → {result} total data rows')
                return max(1, result)

            data_rows = all_rows[1:] if len(all_rows) > 1 else all_rows

            # _apply_content_row_heights will bump every data row to at least
            # _MIN_DATA_ROW_HEIGHT_EMU after filling.  Use that same minimum here so
            # the pre-fill split threshold is consistent with the post-fill actual heights.
            effective_min = self._MIN_DATA_ROW_HEIGHT_EMU

            # Total height assuming template rows are bumped to at least effective_min.
            template_data_h = sum(
                max(r.height or default_row_h, effective_min) for r in data_rows
            )
            template_data_count = len(data_rows)

            # New rows cloned from last template row — also subject to effective_min.
            last_row_h = (data_rows[-1].height or default_row_h) if data_rows else default_row_h
            new_row_h = max(last_row_h, effective_min)

            remaining = max(0, avail_for_data - template_data_h)
            max_new = int(remaining / new_row_h) if new_row_h else 0
            result = template_data_count + max_new

            print(f'  - max_rows_for_shape: avail={avail_height/emu:.2f}" '
                  f'header={header_h/emu:.2f}" '
                  f'tmpl_data={template_data_h/emu:.2f}" ({template_data_count} rows) '
                  f'new_row={new_row_h/emu:.2f}" → {result} total data rows')
            return max(1, result)
        except Exception as e:
            fallback = max(1, self._SVN_SAFE_BOTTOM_EMU // self._MIN_ROW_HEIGHT_EMU)
            print(f'  - _max_rows_for_shape error ({e}); fallback={fallback}')
            return fallback

    def _normalize_table_col_widths(self, table, min_col_width_emu: int = 914400):
        """Ensure no column is narrower than *min_col_width_emu* (default 1 inch).

        When a column is below the minimum, it is expanded to the minimum and the
        surplus is subtracted proportionally from the wider columns, so the total
        table width stays the same.
        """
        try:
            cols = list(table.columns)
            if not cols:
                return
            widths = [c.width for c in cols]
            total_width = sum(widths)

            needs_expand = [w < min_col_width_emu for w in widths]
            if not any(needs_expand):
                return  # all columns already wide enough

            new_widths = list(widths)
            for i, below in enumerate(needs_expand):
                if below:
                    new_widths[i] = min_col_width_emu

            # Amount already locked to narrow columns
            locked_total = sum(w for i, w in enumerate(new_widths) if needs_expand[i])
            remaining = total_width - locked_total
            free_total = sum(w for i, w in enumerate(widths) if not needs_expand[i])

            if free_total > 0 and remaining > 0:
                scale = remaining / free_total
                for i, below in enumerate(needs_expand):
                    if not below:
                        new_widths[i] = max(min_col_width_emu, int(widths[i] * scale))

            # Fix rounding drift so total is exactly preserved
            diff = total_width - sum(new_widths)
            if diff != 0:
                # Apply diff to the widest free column
                widest = max((i for i, b in enumerate(needs_expand) if not b),
                             key=lambda i: new_widths[i], default=0)
                new_widths[widest] += diff

            for col, w in zip(cols, new_widths):
                col.width = w

            print(f'  - Normalized col widths: {[round(w / 914400, 2) for w in new_widths]}" each')
        except Exception as e:
            print(f'  - Error normalizing table col widths: {e}')

    def _fit_table_to_bounds(self, shape):
        """Constrain the table inside the slide's safe content area.

        After rows are added/cloned the table can grow taller than the area
        above the footer.  This method:
          1. Measures total row height after fill.
          2. If it fits → updates shape.height to the actual content height
             (removes dead space / avoids shape box drifting below table).
          3. If it overflows → scales row heights down toward each row's real
             content-required floor, but never below it. A PPTX row's `h` is
             only a *minimum* — PowerPoint/LibreOffice auto-grow a row to fit
             its real text regardless of what we declare, so compressing past
             the content-required height doesn't actually shrink anything; it
             just makes this tool believe the table fits while the real
             renderer still overflows it into the footer.
        """
        try:
            if not shape.has_table:
                return
            table = shape.table
            rows = list(table.rows)
            if not rows:
                return

            total_h = sum(r.height or 0 for r in rows)
            avail_h = max(0, self._SVN_SAFE_BOTTOM_EMU - shape.top)

            if total_h <= avail_h:
                shape.height = total_h   # sync shape box to actual content
                return

            # Need to compress — but only the slack above each row's real
            # content-required height; never below what the text itself needs.
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            col_widths_inch = [max(c.width / self._EMU_PER_INCH, 0.5) for c in table.columns]
            tr_elements = table._tbl.findall('a:tr', namespaces=ns)
            floors = []
            for idx, tr in enumerate(tr_elements):
                if idx == 0:
                    floors.append(self._MIN_ROW_HEIGHT_EMU)
                    continue
                cells = [
                    ''.join(t.text or '' for t in tc.findall('.//a:t', namespaces=ns)).strip()
                    for tc in tr.findall('a:tc', namespaces=ns)
                ]
                floors.append(max(
                    self._MIN_ROW_HEIGHT_EMU,
                    self._estimate_row_content_height_emu(cells, col_widths_inch),
                ))

            floor_total = sum(floors)
            if floor_total >= avail_h:
                print(f'  - _fit_table_to_bounds: content floor '
                      f'{floor_total / 914400:.2f}" already ≥ avail '
                      f'{avail_h / 914400:.2f}" — cannot compress further without '
                      f'clipping text; table may overflow the footer')
                return

            slack_total = total_h - floor_total
            scale = (avail_h - floor_total) / slack_total if slack_total else 0
            new_heights = [
                floor + int(((r.height or 0) - floor) * scale)
                for r, floor in zip(rows, floors)
            ]

            # Fix rounding drift: adjust the tallest row
            diff = avail_h - sum(new_heights)
            if diff != 0:
                idx = max(range(len(new_heights)), key=lambda i: new_heights[i])
                new_heights[idx] += diff

            for row, h in zip(rows, new_heights):
                row.height = h
            shape.height = sum(r.height for r in rows)
            print(f'  - Fit table to safe area: '
                  f'{total_h / 914400:.2f}" → {shape.height / 914400:.2f}" '
                  f'(safe bottom {self._SVN_SAFE_BOTTOM_EMU / 914400:.2f}")')
        except Exception as e:
            print(f'  - _fit_table_to_bounds error: {e}')

    def _find_shape(self, slide, target):
        """Find shape by index or name.

        When only ``shape_name`` is provided (no ``shape_index``), the search
        is performed recursively through GROUP shapes so that shapes nested
        inside groups can be targeted by name.
        """
        try:
            # Try by index first (faster, top-level only)
            if target.shape_index is not None and target.shape_index < len(slide.shapes):
                shape = slide.shapes[target.shape_index]
                print(f'  - Found shape by index {target.shape_index}: {shape.name}')
                return shape

            # Fallback to recursive name search (traverses group children)
            if target.shape_name:
                found = self._find_shape_recursive(slide.shapes, target.shape_name)
                if found:
                    print(f'  - Found shape by name (recursive): {target.shape_name}')
                    return found
        except Exception as e:
            print(f'  - Error finding shape: {e}')

        return None

    def _find_shape_recursive(self, shapes, name: str):
        """Recursively search for a shape by name, traversing GROUP children.

        Args:
            shapes: A shape collection (``slide.shapes`` or ``group_shape.shapes``).
            name: The exact shape name to search for.

        Returns:
            The matching shape, or ``None`` if not found.
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            for shape in shapes:
                if shape.name == name:
                    return shape
                # Recurse into group shapes
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        found = self._find_shape_recursive(shape.shapes, name)
                        if found:
                            return found
                except AttributeError:
                    pass
        except Exception as e:
            print(f'  - _find_shape_recursive error: {e}')
        return None

    def _fill_slide_default(self, slide, slide_content: SlideContent):
        """Default filling strategy (original logic)"""

        # Update title
        self._update_slide_title(slide, slide_content.title)

        # Get available text containers
        placeholders = self._get_text_placeholders(slide)
        text_shapes = self._get_text_shapes(slide)

        # Fill content items
        placeholder_idx = 0
        for content_item in slide_content.content:
            if content_item.type == ContentType.TEXT:
                placeholder_idx = self._fill_text_content(
                    slide, content_item, placeholders, text_shapes, placeholder_idx,
                )

            elif content_item.type == ContentType.LIST:
                placeholder_idx = self._fill_list_content(
                    slide, content_item, placeholders, placeholder_idx,
                )

            elif content_item.type == ContentType.TABLE:
                self._add_table(slide, content_item, Inches(1.0), Inches(2.5), Inches(8.0))

            elif content_item.type == ContentType.IMAGE:
                self._add_image(slide, content_item, Inches(1.0), Inches(2.5), Inches(8.0))

    def _get_text_placeholders(self, slide):
        """Get all text placeholders in slide (excluding title)"""
        return [
            s for s in slide.shapes
            if s.is_placeholder and s.has_text_frame and s != slide.shapes.title
        ]

    def _get_text_shapes(self, slide):
        """Get all text shapes in slide (excluding placeholders and title)"""
        return [
            s for s in slide.shapes
            if s.has_text_frame and not s.is_placeholder and s != slide.shapes.title
        ]

    def _fill_text_content(self, slide, content_item, placeholders, text_shapes, placeholder_idx):
        """Fill text content into placeholders or text shapes"""
        if placeholder_idx < len(placeholders):
            placeholders[placeholder_idx].text_frame.text = content_item.data
            self._apply_text_font(placeholders[placeholder_idx].text_frame)
            print(f'  - Filled placeholder {placeholder_idx}: {content_item.data[:50]}...')
            return placeholder_idx + 1

        elif text_shapes:
            text_shapes[0].text_frame.text = content_item.data
            self._apply_text_font(text_shapes[0].text_frame)
            text_shapes.pop(0)
            return placeholder_idx

        else:
            # Add new textbox if no placeholders left
            self._add_content_at_position(slide, content_item, Inches(1.0), Inches(2.0), Inches(8.0))
            return placeholder_idx

    def _fill_list_content(self, slide, content_item, placeholders, placeholder_idx):
        """Fill list content into placeholder"""
        if placeholder_idx < len(placeholders):
            text_frame = placeholders[placeholder_idx].text_frame
            text_frame.clear()
            for i, item in enumerate(content_item.data):
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                p.text = f'• {item}'
                p.level = 0
            self._apply_text_font(text_frame)
            print(f'  - Filled list with {len(content_item.data)} items')
            return placeholder_idx + 1
        else:
            # Add new list
            top = Inches(2.5 + (placeholder_idx * 0.5))
            self._add_content_at_position(slide, content_item, Inches(1.0), top, Inches(8.0))
            return placeholder_idx

    def _add_content_at_position(self, slide, content_item: ContentItem, left, top, width):
        """Add content item at specific position"""
        if content_item.type == ContentType.TEXT:
            height = Inches(0.5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text_frame.text = content_item.data
            self._apply_text_font(textbox.text_frame)

        elif content_item.type == ContentType.LIST:
            height = Inches(0.3) * len(content_item.data)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame

            for i, item in enumerate(content_item.data):
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                p.text = f'• {item}'
            self._apply_text_font(text_frame)

    def _add_image(self, slide, content_item: ContentItem, left, top, width):
        """Add image to slide"""
        try:
            image_data = content_item.data
            image_path = self._resolve_image_path(image_data)

            if Path(image_path).exists():
                slide.shapes.add_picture(str(image_path), left, top, width=width)
                self._add_image_caption(slide, image_data, left, top, width)
            else:
                print(f'Image not found: {image_path}')
                self._add_placeholder_text(slide, f'[Image: {image_path}]', left, top, width)

        except Exception as e:
            print(f'Error adding image: {e}')

    def _resolve_image_path(self, image_data) -> str:
        """Resolve image path from data"""
        image_path = image_data.get('path') if isinstance(image_data, dict) else image_data

        if image_path and not Path(str(image_path)).is_absolute():
            # Try relative to CWD first
            rel_path = Path(str(image_path))
            if not rel_path.exists():
                # Try outputs directory
                image_path = str(Path('outputs') / str(image_path))

        return str(image_path) if image_path else ''

    def _add_image_caption(self, slide, image_data, left, top, width):
        """Add caption below image if exists"""
        if isinstance(image_data, dict) and image_data.get('caption'):
            caption_top = top + Inches(3.0)
            caption_box = slide.shapes.add_textbox(left, caption_top, width, Inches(0.3))
            caption_box.text_frame.text = image_data['caption']
            for paragraph in caption_box.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Noto Sans JP'
                    run.font.size = Pt(10)
                    run.font.italic = True

    @staticmethod
    def _set_text_preserving_format(text_frame, value: str) -> None:
        """Set a text frame's text, reusing the first existing run instead of
        `text_frame.text = value`. That setter replaces every paragraph/run
        with a single freshly-styled one, silently discarding any
        template-authored run formatting (e.g. this template's "after" title
        cells are colored/bold via their own run, not the placeholder's
        paragraph/shape defaults — resetting the run drops that styling and
        the text renders in whatever default color the theme falls back to,
        which can be invisible against the slide's own background)."""
        paragraphs = text_frame.paragraphs
        if paragraphs and paragraphs[0].runs:
            first_para = paragraphs[0]
            first_para.runs[0].text = value
            for extra_run in list(first_para.runs[1:]):
                extra_run._r.getparent().remove(extra_run._r)
            for extra_para in list(paragraphs[1:]):
                extra_para._p.getparent().remove(extra_para._p)
        else:
            text_frame.text = value

    def _apply_text_font(self, text_frame, font_name: str = 'Noto Sans JP', font_size: int = 10):
        """Apply font and size to all text in a text frame"""
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = Pt(font_size)
        except Exception as e:
            print(f'  - Error applying font to text: {e}')

    def _apply_line_spacing(self, text_frame, spacing_pt):
        """Apply line spacing (space after) to all paragraphs in a text frame
        
        Args:
            text_frame: The text frame to apply spacing to
            spacing_pt: Spacing in points (e.g., Pt(10) for 10pt)
        """
        try:
            for paragraph in text_frame.paragraphs:
                if paragraph.text.strip():  # skip empty paragraphs (from blank lines)
                    paragraph.space_after = spacing_pt
        except Exception as e:
            print(f'  - Error applying line spacing: {e}')

    def _estimate_max_chars(self, shape, font_pt: int = 8) -> int:
        """Estimate max characters that fit in a shape based on its dimensions.

        Uses CJK character width at the given font size to compute chars-per-line
        and shape height to compute how many lines fit.  Returns a floor of 8 so
        at least a short label is always written.
        """
        try:
            emu = self._EMU_PER_INCH
            w_inch = (shape.width or 0) / emu
            h_inch = (shape.height or 0) / emu
            cjk_w = font_pt / 72          # ~0.111" per CJK char at 8pt
            line_h = font_pt * 1.35 / 72  # ~0.150" line height with spacing
            chars_per_line = max(5, int(w_inch / cjk_w))
            lines = max(1, int(h_inch / line_h))
            return max(8, chars_per_line * lines)
        except Exception:
            return 65

    def _estimate_extra_wrap_height(self, shape, text: str, font_pt: int = 10) -> int:
        """EMU height needed for `text` to word-wrap in `shape`'s current
        width, beyond what a single line already occupies. Used to grow a
        title box to fit its full text instead of truncating it with an
        ellipsis — a short heading cut off mid-word reads as meaningless."""
        try:
            emu = self._EMU_PER_INCH
            w_inch = (shape.width or 0) / emu
            cjk_w = font_pt / 72
            chars_per_line = max(5, int(w_inch / cjk_w))
            lines_needed = -(-len(text) // chars_per_line) if text else 1  # ceil
            extra_lines = max(0, lines_needed - 1)
            line_h_inch = font_pt * 1.35 / 72
            return int(extra_lines * line_h_inch * emu)
        except Exception:
            return 0

    @staticmethod
    def _set_vertical_anchor_top(text_frame) -> None:
        """Anchor text to the top of the shape so it doesn't float to center."""
        try:
            from pptx.enum.text import MSO_ANCHOR
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
        except Exception:
            pass

    @staticmethod
    def _truncate_at_sentence(text: str, max_chars: int = 120) -> str:
        """Truncate text to fit within max_chars, preferring sentence boundaries.

        Tries Japanese sentence ends (。) then newlines then hard cut.
        Returns the original string unchanged when it is already short enough.
        """
        if len(text) <= max_chars:
            return text
        truncated = text[:max_chars]
        for sep in ['。', '\n', '．', '.']:
            idx = truncated.rfind(sep)
            if idx > max_chars // 2:
                return truncated[:idx + 1].rstrip()
        return truncated.rstrip() + '…'

    def _apply_auto_fit(self, text_frame):
        """Enable normAutofit so text shrinks to fit the shape instead of overflowing."""
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception as e:
            print(f'  - Could not apply auto-fit: {e}')

    def _apply_table_cell_font(self, cell, font_name: str = 'Noto Sans JP', font_size: int = 8):
        """Apply font and size to all text in a table cell"""
        try:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = Pt(font_size)
        except Exception as e:
            print(f'  - Error applying font to table cell: {e}')

    def _add_placeholder_text(self, slide, text: str, left, top, width):
        """Add placeholder text box"""
        textbox = slide.shapes.add_textbox(left, top, width, Inches(1))
        textbox.text_frame.text = text
        self._apply_text_font(textbox.text_frame)

    def _add_table(self, slide, content_item: ContentItem, left, top, width):
        """Add table to slide"""
        try:
            table_data = content_item.data
            headers = table_data.get('headers', [])
            rows = table_data.get('rows', [])

            if not headers or not rows:
                print('Table has no headers or rows')
                return

            # Create table
            row_count = len(rows) + 1  # +1 for header
            col_count = len(headers)
            height = Inches(0.3) * row_count

            table_shape = slide.shapes.add_table(row_count, col_count, left, top, width, height)
            table = table_shape.table

            # Fill headers
            self._fill_table_headers(table, headers)

            # Fill data rows
            self._fill_table_rows(table, rows)

        except Exception as e:
            print(f'Error adding table: {e}')

    def _fill_table_headers(self, table, headers: List[str]):
        """Fill table header row"""
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            # Apply font and make header bold
            self._apply_table_cell_font(cell)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

    def _fill_table_rows(self, table, rows: List[List[str]]):
        """Fill table data rows"""
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)  # +1 to skip header
                cell.text = str(cell_data)
                self._apply_table_cell_font(cell)

    # ==================================================================
    # Role-based rendering (new pipeline: profile.md -> ProjectProfile -> PPTX)
    # ==================================================================

    def render_from_profile(
        self,
        profile: ProjectProfile,
        template: str,
        output_name: str,
        output_dir: Optional[str] = None,
        extra_slides: Optional[List[dict]] = None,
    ) -> Path:
        """Fill template using a ProjectProfile + role-based slide configs.

        For each configured slide:
          - Resolve each role's data via source_path
          - Discover shapes via ShapeDiscovery
          - Fill shapes (text / table / image / compound pairs)
          - Clear excess shapes
        """
        print(f'Rendering from profile: {output_name}')
        self._output_dir = self._resolve_output_dir(output_dir)
        self._overflow_extras = {}
        self._detect_template_config(template)
        prs = self._load_template(template)
        # Compute safe content area from actual slide height so tables never reach the footer.
        # Footer band occupies roughly the bottom 13% of the slide; 0.87 leaves ~0.36" margin
        # above a footer at 93% (observed: 5.254" on a 5.625" slide → 93.4%).
        if prs.slide_height:
            self._SVN_SAFE_BOTTOM_EMU = int(prs.slide_height * 0.87)
            print(f'  Safe content bottom: {self._SVN_SAFE_BOTTOM_EMU / self._EMU_PER_INCH:.3f}" '
                  f'(slide is {prs.slide_height / self._EMU_PER_INCH:.3f}" tall)')
        if prs.slide_width:
            self._slide_width_emu = prs.slide_width
        # Captured before pruning: pruning deletes the trailing blank
        # end-card slides that make the safest fallback canvas for extra
        # slides, but the SlideLayout object stays valid (it's a separate
        # XML part from the slide instance) so it can still be reused after.
        blank_canvas_layout = self._find_blank_canvas_layout(prs)
        spans = self._fill_slides_from_profile(prs, profile)
        last_positions = self._prune_and_reorder(prs, profile, spans)
        all_extra_slides = list(extra_slides or []) + self._build_overflow_extra_slides()
        if all_extra_slides:
            self._render_extra_slides(prs, all_extra_slides, last_positions, blank_canvas_layout)
        return self._save_presentation(prs, output_name, output_dir)

    def _build_overflow_extra_slides(self) -> List[dict]:
        """Turn queued per-item text overflow (`_queue_overflow_extra`) into
        appendix extra-slide entries, one per (section, role) group, anchored
        right after that section so the full text lands near where it was cut."""
        entries = []
        for (section_key, role_name), texts in self._overflow_extras.items():
            label = role_name.replace('_', ' ').title()
            entries.append({
                'layout': 'bullets',
                'title': f'{label} — Additional Detail',
                'bullets': texts,
                'anchor_section': section_key,
            })
        return entries

    def _find_blank_canvas_layout(self, prs: Presentation):
        """Find a clean, minimal-chrome layout to fall back on for extra slides.

        Scans from the end of the deck since the SVN template's trailing
        slides are blank end-cards with no decorative content-slide styling
        (unlike e.g. project_background's two-tone split-panel master).
        """
        slide_area = (prs.slide_width or 0) * (prs.slide_height or 0)
        for slide in reversed(prs.slides):
            if self._is_clean_master(slide.slide_layout.slide_master, slide_area):
                return slide.slide_layout
        return None

    def _fill_slides_from_profile(
        self, prs: Presentation, profile: ProjectProfile,
    ) -> Dict[int, Tuple[int, int]]:
        """Fill configured slides; return map {template_slide_num: (start_idx, end_idx)}.

        ``start_idx`` is the slide's current (0-based) index at fill time;
        ``end_idx`` is the current index of its last overflow-continuation
        slide (``end_idx == start_idx`` when the slide didn't overflow). The
        prune+reorder pass uses this block to keep overflow continuations
        contiguous with their parent.
        """
        spans: Dict[int, Tuple[int, int]] = {}
        if not self.template_config:
            print('No template config — cannot render from profile')
            return spans

        configs = self.template_config.slide_role_configs
        sorted_nums = sorted(configs.keys())
        slide_offset = 0
        discovery = ShapeDiscovery()

        for slide_num in sorted_nums:
            if not self.template_config.should_fill_slide(slide_num):
                print(f'Slide {slide_num} is protected — skipping')
                continue
            adj_idx = slide_num - 1 + slide_offset
            if adj_idx >= len(prs.slides):
                print(f'Slide {slide_num} not found in template (have {len(prs.slides)})')
                continue
            slide = prs.slides[adj_idx]
            role_config = configs[slide_num]
            print(f'Filling slide {slide_num} (template index {adj_idx}): '
                  f'{len(role_config.roles)} role(s)')
            extra = self._fill_slide_roles(prs, slide, role_config, profile, adj_idx, discovery)
            # Span covers the parent slide + any overflow continuation slides.
            spans[slide_num] = (adj_idx, adj_idx + extra)
            slide_offset += extra
        return spans

    # ------------------------------------------------------------------
    # Prune + reorder pass — drop template slides not backed by content,
    # keep chapter dividers only when their chapter has content, and order
    # the survivors by `profile.section_order`.
    # ------------------------------------------------------------------

    def _prune_and_reorder(
        self, prs: Presentation, profile: ProjectProfile, spans: Dict[int, Tuple[int, int]],
    ) -> Dict[int, int]:
        """Delete unwanted template slides and reorder survivors.

        Returns ``new_last_positions``: {template_slide_num: new_0based_index}
        for every configured slide that survived — the fresh position map
        ``_render_extra_slides`` (and its chrome/canvas pickers) must use,
        since positions computed before this pass are stale afterward.

        No-ops (returns positions unchanged) when there's no recognized
        template config — this pass hardcodes SVN's slide-picking semantics
        (``CHAPTER_DIVIDERS`` etc.), so it must never run against a template
        `_detect_template_config` didn't recognize, or it would destructively
        delete slides using slide numbers that mean nothing for that deck.
        """
        if not self.template_config:
            return {n: e for n, (s, e) in spans.items()}

        from lib.templates.svn import ALWAYS_KEEP_SLIDES, CHAPTER_DIVIDERS, SECTION_TO_SLIDES

        pml_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        r_id_attr = f'{{{r_ns}}}id'

        # Reverse lookup: content slide number -> the chapter divider that introduces it.
        slide_to_divider = {c: d for d, cs in CHAPTER_DIVIDERS.items() for c in cs}

        def _current_index(orig: int) -> int:
            """Current (post-fill, pre-prune) 0-based index of an un-configured
            slide number (cover/agenda/divider), accounting for overflow slides
            inserted before it by earlier configured sections."""
            return (orig - 1) + sum(e - s for m, (s, e) in spans.items() if m < orig)

        # --- (a)+(b) build wanted_orig: always-kept + content sections' slides,
        #     strictly in `profile.section_order` (i.e. markdown heading order).
        #     The user owns the presentation order they write in the markdown —
        #     sections are never reshuffled into template/canonical order. Each
        #     chapter divider is inserted once, right before the first sibling
        #     of that chapter to appear in the markdown; later-appearing
        #     siblings of the same chapter are emitted in place (no duplicate
        #     divider) since the template only has one divider slide per chapter.
        order_source = profile.section_order or list(SECTION_TO_SLIDES.keys())

        wanted_orig: list = list(ALWAYS_KEEP_SLIDES)
        emitted_keys: set = set()
        dividers_emitted: set = set()
        for key in order_source:
            if key in emitted_keys:
                continue
            slide_nums = SECTION_TO_SLIDES.get(key)
            if not slide_nums:
                continue
            if not section_has_content(profile, key):
                continue
            divider = slide_to_divider.get(slide_nums[0])
            if divider is not None and divider not in dividers_emitted:
                wanted_orig.append(divider)
                dividers_emitted.add(divider)
            wanted_orig.extend(slide_nums)
            emitted_keys.add(key)

        # --- (c) expand each wanted original slide number to its current-index
        #     block: configured slides expand to their full overflow block;
        #     everything else (cover/agenda/dividers) resolves to one index.
        wanted_indices_ordered: list = []
        seen_idx = set()
        for orig in wanted_orig:
            if orig in spans:
                start, end = spans[orig]
                idx_block = range(start, end + 1)
            else:
                idx_block = [_current_index(orig)]
            for idx in idx_block:
                if idx in seen_idx:
                    print(f'  - _prune_and_reorder: duplicate index {idx} '
                          f'(slide {orig}) — skipping')
                    continue
                wanted_indices_ordered.append(idx)
                seen_idx.add(idx)

        keep_set = set(wanted_indices_ordered)

        # --- (d) sldIdLst surgery: delete unwanted slides (drop_rel so no
        #     orphaned parts remain), reorder survivors to match final order.
        sldIdLst = prs.part._element.find(f'{{{pml_ns}}}sldIdLst')
        if sldIdLst is None:
            print('  - _prune_and_reorder: no <p:sldIdLst> found — skipping prune')
            return {n: e for n, (s, e) in spans.items()}

        snapshot = list(sldIdLst)
        final_elems = [snapshot[i] for i in wanted_indices_ordered]

        for child in list(sldIdLst):
            sldIdLst.remove(child)
        for elem in final_elems:
            sldIdLst.append(elem)

        dropped = 0
        for i, elem in enumerate(snapshot):
            if i in keep_set:
                continue
            rid = elem.get(r_id_attr)
            if rid:
                try:
                    prs.part.drop_rel(rid)
                    dropped += 1
                except KeyError:
                    pass

        print(f'  - _prune_and_reorder: kept {len(final_elems)}/{len(snapshot)} slides '
              f'(dropped {dropped} relationship(s)); kept template slides in order: '
              f'{wanted_orig}')

        # --- (e) fresh positions for surviving configured slides.
        new_last_positions: Dict[int, int] = {}
        for n, (s, e) in spans.items():
            if e in keep_set:
                new_last_positions[n] = wanted_indices_ordered.index(e)
        return new_last_positions

    def _fill_slide_roles(
        self, prs: Presentation, slide, role_config: SlideRoleConfig,
        profile: ProjectProfile, slide_index: int, discovery: ShapeDiscovery,
    ) -> int:
        """Fill all roles on a slide. Returns extra slides inserted (overflow continuations).

        TABLE roles keep their existing independent overflow — no slide in this
        template mixes a table with other overflow-prone roles. TEXT-kind roles
        (collections / compound title+body pairs) instead share ONE set of
        continuation slides per base slide, sized to the largest role's chunk
        count — so multiple overflowing roles on the same slide (e.g. business
        process's `category_labels` + `after_blocks`) land on the same
        continuation pages instead of each duplicating the slide independently,
        which would corrupt non-overflowing sibling roles on the resulting
        chain of duplicates.
        """
        if role_config.slide_number == 11:
            split_extra = self._maybe_split_business_process_slide(
                prs, slide, role_config, profile, slide_index, discovery)
            if split_extra is not None:
                return split_extra

        extra_total = 0
        text_chunks = []  # [(role, [chunk0, chunk1, ...]), ...]
        for role in role_config.roles:
            data = self._resolve_source_path(profile, role.source_path)
            shapes = discovery.discover(slide, role)
            if role.kind == ContentKind.TABLE:
                extra_total += self._fill_table_role(prs, slide, slide_index + extra_total,
                                                    shapes, data, role)
            elif role.kind == ContentKind.IMAGE:
                self._fill_image_role(slide, shapes, data, role)
            elif role.sub_roles:  # compound (title+body pairs)
                items = data if isinstance(data, list) else []
                chunks = self._chunk_list(items, len(shapes))
                self._fill_compound_pairs(shapes, chunks[0], role, slide=slide)
                text_chunks.append((role, chunks))
            else:  # simple text / text collection
                items = self._coerce_text_items(data)
                chunks = self._chunk_list(items, len(shapes))
                self._fill_text_shapes(shapes, chunks[0], role, slide=slide)
                text_chunks.append((role, chunks))

        max_chunks = max((len(chunks) for _, chunks in text_chunks), default=1)
        if max_chunks > 1:
            overflowing = [r.name for r, chunks in text_chunks if len(chunks) > 1]
            print(f'  - Text overflow: {max_chunks} slide(s) needed (role(s): {overflowing})')
            current_index = slide_index + extra_total
            for k in range(1, max_chunks):
                new_slide = self._duplicate_slide(prs, current_index)
                extra_total += 1
                current_index += 1
                for role, chunks in text_chunks:
                    # Roles with fewer chunks than max_chunks have nothing new to
                    # show on this continuation — clear them rather than repeating
                    # their last chunk verbatim, which reads as a mistake (the
                    # same items appearing twice) rather than intentional content.
                    chunk = chunks[k] if k < len(chunks) else []
                    new_shapes = ShapeDiscovery().discover(new_slide, role)
                    if role.sub_roles:
                        self._fill_compound_pairs(new_shapes, chunk, role, slide=new_slide)
                    else:
                        self._fill_text_shapes(new_shapes, chunk, role, slide=new_slide)
                print(f'  - Filled continuation {k}/{max_chunks - 1}')
        return extra_total

    @staticmethod
    def _chunk_list(items: list, capacity: int) -> list:
        """Split items into chunks of `capacity` size each (always >=1 chunk)."""
        if capacity <= 0:
            return [items]
        return [items[i:i + capacity] for i in range(0, len(items), capacity)] or [[]]

    def _maybe_split_business_process_slide(
        self, prs: Presentation, slide, role_config: SlideRoleConfig,
        profile: ProjectProfile, slide_index: int, discovery: ShapeDiscovery,
    ) -> Optional[int]:
        """Bespoke fallback for template slide 11 (business process, horizontal
        view): duplicating the whole slide for `after_blocks` overflow would
        repeat the unrelated `before_steps` row pointlessly on the
        continuation. Instead split into a before-only slide (trimmed to the
        real step count) and an after-only slide with the overflow items laid
        out as a second row, reusing the existing column shapes — closer to
        what a human editor would do by hand.

        Returns None (caller falls back to the generic per-role handling)
        when `after_blocks` isn't overflowing, or when it overflows by more
        than a second row can hold (2 rows of `capacity` each) — the generic
        duplicate-and-paginate path is the safe fallback for that case.
        """
        before_role = next((r for r in role_config.roles if r.name == 'before_steps'), None)
        after_role = next((r for r in role_config.roles if r.name == 'after_blocks'), None)
        if before_role is None or after_role is None:
            return None

        after_items = self._resolve_source_path(profile, after_role.source_path)
        after_items = after_items if isinstance(after_items, list) else []
        row1_pairs = discovery.discover(slide, after_role)
        capacity = len(row1_pairs)
        if capacity == 0 or not (capacity < len(after_items) <= capacity * 2):
            return None

        from lib.templates.svn import BUSINESS_PROCESS_SLIDE11_LAYOUT as L

        def shape_at(target_slide, idx):
            return ShapeDiscovery._safe_get_by_index(target_slide, idx)

        print(f'  - Role "after_blocks": {len(after_items)} items > {capacity} slots — '
              f'splitting into before-only + after-only (2 rows) slides')

        before_items = self._coerce_text_items(
            self._resolve_source_path(profile, before_role.source_path))
        before_boxes = discovery.discover(slide, before_role)
        n_before = len(before_items)
        n_row2 = len(after_items) - capacity

        # Resolve every shape reference we'll need — on BOTH the original
        # slide and its as-yet-uncreated duplicate — before any deletion.
        # Removing a shape shifts the live index of every shape after it, so
        # re-looking-up by index mid-deletion would silently grab the wrong
        # shape; every index lookup below happens against untouched slides.
        before_connectors = [shape_at(slide, i) for i in L['before_connectors']]
        before_bg = shape_at(slide, L['before_bg'])
        before_bg_top = before_bg.top
        after_bg = shape_at(slide, L['after_bg'])
        after_bg_top = after_bg.top
        after_label_top = shape_at(slide, L['after_label']).top
        after_card_bg_top = shape_at(slide, L['after_card_bg'][0]).top
        after_title_top = shape_at(slide, L['after_title_boxes'][0]).top
        after_body_top = shape_at(slide, L['after_body_boxes'][0]).top
        # Connectors sit left-to-right between consecutive columns; the raw
        # index order in the template isn't left-to-right, so sort by position.
        after_connectors_sorted = sorted(
            (shape_at(slide, i) for i in L['after_connectors']), key=lambda s: s.left)
        after_connector_top = after_connectors_sorted[0].top if after_connectors_sorted else after_body_top
        row2_label_top = before_bg_top + (after_label_top - after_bg_top)
        row2_card_bg_top = before_bg_top + (after_card_bg_top - after_bg_top)
        row2_title_top = before_bg_top + (after_title_top - after_bg_top)
        row2_body_top = before_bg_top + (after_body_top - after_bg_top)
        row2_connector_top = before_bg_top + (after_connector_top - after_bg_top)
        after_side_shapes = [shape_at(slide, idx) for idx in [
            L['after_bg'], L['after_label'], L['row_arrow'], *L['after_card_bg'],
            *L['after_title_boxes'], *L['after_body_boxes'], *L['after_connectors'],
        ]]

        # 1. Duplicate the pristine slide first — it becomes the after-only slide.
        after_slide = self._duplicate_slide(prs, slide_index)

        # Resolve the duplicate's own shape references too, before touching it.
        # `row_arrow` is kept (not deleted) on this slide: originally it showed
        # before-state flowing into after-state, and with 2 after-rows stacked
        # it now reads just as naturally as "row 2 flows into row 1" — same
        # gap between panels, no repositioning needed.
        before_side_shapes_on_dup = [shape_at(after_slide, idx) for idx in [
            L['before_label'], *L['before_boxes'], *L['before_connectors'],
        ]]
        before_bg_on_dup = shape_at(after_slide, L['before_bg'])
        after_label_on_dup = shape_at(after_slide, L['after_label'])
        after_card_bg_on_dup = [shape_at(after_slide, i) for i in L['after_card_bg']]
        after_connectors_on_dup = sorted(
            (shape_at(after_slide, i) for i in L['after_connectors']), key=lambda s: s.left)
        row1_after_pairs = discovery.discover(after_slide, after_role)

        # Drop before-side shapes and replace the (differently-styled)
        # before_bg panel with a clone of after_bg repositioned into its spot
        # — both rows share the same panel style instead of one row looking
        # like a leftover "before" panel. This MUST happen before row 2's
        # content is cloned below: `_clone_shape` appends to the end of the
        # shape tree, so cloning the panel afterwards would stack it on top
        # of (hiding) row 2's title/body/card_bg.
        for shp in before_side_shapes_on_dup:
            self._remove_shape(shp)
        self._remove_shape(before_bg_on_dup)
        self._clone_shape(after_slide, after_bg, top=before_bg_top)

        # Normalize row 1's own body cells to their largest height AND width
        # before filling/cloning — several of this template's 4 body cells are
        # much smaller than the others (as little as 0.13" tall / 1.21" wide),
        # which forces heavy truncation. Widening in place (left unchanged) is
        # safe here — checked against this exact template's column spacing,
        # the widened cells still clear the next column's left edge. Row 2
        # clones inherit the fix since it copies row 1 *after* this runs.
        row1_bodies = [b for _, b in row1_after_pairs if b is not None]
        if row1_bodies:
            max_body_h = max(b.height for b in row1_bodies)
            max_body_w = max(b.width for b in row1_bodies)
            for b in row1_bodies:
                b.height = max_body_h
                b.width = max_body_w

        row2_pairs = []
        for i in range(n_row2):
            title_src, body_src = row1_after_pairs[i]
            # Card backdrop cloned first so it stays behind title/body in z-order,
            # matching row 1's own stacking (card_bg shapes precede title/body in
            # the template's original shape order).
            self._clone_shape(after_slide, after_card_bg_on_dup[i], top=row2_card_bg_top)
            new_title = self._clone_shape(after_slide, title_src, top=row2_title_top)
            new_body = self._clone_shape(after_slide, body_src, top=row2_body_top)
            row2_pairs.append((new_title, new_body))
        for i in range(max(0, n_row2 - 1)):
            self._clone_shape(after_slide, after_connectors_on_dup[i], top=row2_connector_top)

        # 2. Original slide -> before-only: fill + trim + widen before_steps to
        #    fill the freed-up width instead of staying bunched at their
        #    original (8-box) narrow size, drop after-side shapes.
        self._fill_text_shapes(before_boxes, before_items, before_role)
        kept_boxes = before_boxes[:n_before]
        kept_connectors = before_connectors[:max(0, n_before - 1)]
        for shp in before_boxes[n_before:]:
            self._remove_shape(shp)
        for shp in before_connectors[max(0, n_before - 1):]:
            self._remove_shape(shp)
        if kept_boxes:
            span_left = before_boxes[0].left
            span_right = before_boxes[-1].left + before_boxes[-1].width
            gap = before_connectors[0].width if before_connectors else 0
            box_w = int((span_right - span_left - gap * (n_before - 1)) / n_before)
            x = span_left
            for i, box in enumerate(kept_boxes):
                box.left, box.width = int(x), box_w
                x += box_w
                if i < len(kept_connectors):
                    kept_connectors[i].left = int(x)
                    x += gap
        for shp in after_side_shapes:
            self._remove_shape(shp)

        # 3. Duplicated slide -> after-only, 2 rows: fill row 1 (existing
        #    columns) + row 2 (cloned columns from step 1), and move the
        #    single "after" label above both rows.
        if after_label_on_dup is not None:
            after_label_on_dup.top = row2_label_top
        self._fill_compound_pairs(row1_after_pairs, after_items[:capacity], after_role)
        self._fill_compound_pairs(row2_pairs, after_items[capacity:], after_role)

        return 1

    # ------------------------------------------------------------------
    # Source-path resolution: dot notation + optional [start:end] slice
    # ------------------------------------------------------------------

    def _resolve_source_path(self, profile: ProjectProfile, path: str):
        """Resolve `a.b.c` or `a[0:2]` or `a.b[2:4]` into a value on the profile.

        Returns None if any segment is missing.
        """
        if not path:
            return None
        # Split off trailing slice like benefits[0:2]
        slice_spec = None
        if path.endswith(']') and '[' in path:
            base, slc = path.rsplit('[', 1)
            slice_spec = slc[:-1]  # strip ']'
            path = base
        cur = profile
        for part in path.split('.'):
            if cur is None:
                return None
            cur = getattr(cur, part, None)
        if slice_spec is not None and isinstance(cur, list):
            try:
                a, b = slice_spec.split(':')
                start = int(a) if a else None
                end = int(b) if b else None
                return cur[start:end]
            except Exception:
                return cur
        return cur

    # ------------------------------------------------------------------
    # Role fillers
    # ------------------------------------------------------------------

    def _fill_text_shapes(self, shapes: list, items: list, role: ShapeRole, slide=None):
        """Fill one page's worth of text into shapes; clears unused slots.

        An item too long for its fixed box is truncated in place at a
        sentence boundary — same policy `_fill_compound_pairs` uses for body
        text. These boxes (e.g. business-process step/category labels) sit
        in diagram layouts next to sibling shapes/connectors, so they can't
        safely grow like a title/body pair can; the full text is queued for
        an auto-generated appendix extra-slide instead (`_queue_overflow_extra`)
        so nothing is lost, and the slide itself is never duplicated just to
        make room for one long item (that would repeat every OTHER role on
        the same slide for no reason).

        When `slide` is given, an unused slot's decoration (icon/background —
        not tracked by the role, so blanking text alone would leave an
        empty-looking box) is removed via `_remove_decorations_around`.
        Callers with their own bespoke unused-slot handling (e.g.
        `_maybe_split_business_process_slide`) omit `slide` to skip this.
        """
        for i, shape in enumerate(shapes):
            if shape is None:
                continue
            if i >= len(items):
                if shape.has_text_frame:
                    self._set_text_preserving_format(shape.text_frame, '')
                if slide is not None:
                    self._remove_decorations_around(slide, [shape])
                continue
            value = items[i]
            if shape.has_text_frame:
                shape.text_frame.word_wrap = True
                max_chars = self._estimate_max_chars(shape, font_pt=10)
                if max_chars <= self._MIN_MEANINGFUL_TRUNCATION_CHARS:
                    display_value, truncated = value, False
                else:
                    display_value = self._truncate_at_sentence(value, max_chars=max_chars)
                    truncated = display_value != value
                    if truncated:
                        self._queue_overflow_extra(role, value)
                self._set_text_preserving_format(shape.text_frame, display_value)
                self._apply_text_font(shape.text_frame)
                # Matches _fill_compound_pairs's body handling: once text is
                # truncated to the estimated capacity, skip auto-fit — on
                # some renderers (LibreOffice) normAutofit's fit calculation
                # can still shrink already-truncated text to near-invisible
                # sizes. Untruncated text (either it fit already, or the box
                # was too small to truncate meaningfully) keeps auto-fit as
                # a safety net.
                if not truncated:
                    self._apply_auto_fit(shape.text_frame)
                if role.name == 'current_issues' or role.name == 'objectives':
                    self._apply_line_spacing(shape.text_frame, Pt(10))

    def _queue_overflow_extra(self, role: ShapeRole, full_text: str) -> None:
        """Record an item's full (untruncated) text for an appendix extra-slide.

        Grouped by (section, role name) so e.g. multiple long `before_steps`
        items land on one shared appendix slide rather than one each.
        Rendered by `render_from_profile` after the main fill+prune pass —
        see `_overflow_extras`.
        """
        section_key = (role.source_path or '').split('.')[0].split('[')[0] or role.name
        bucket = self._overflow_extras.setdefault((section_key, role.name), [])
        if full_text not in bucket:
            bucket.append(full_text)
        print(f"  - Role '{role.name}': item too long for its box — truncated in place, "
              f'full text queued for an appendix slide')

    def _fill_compound_pairs(self, pairs: list, items: list, role: ShapeRole, slide=None):
        """Fill a single slide's worth of (title, body) pairs; clears unused slots.

        When `slide` is given, an unused pair's decoration (card background,
        icon badge, icon picture — not tracked by the role, so blanking
        title/body text alone would leave an empty-looking card) is removed
        via `_remove_decorations_around`. Callers with their own bespoke
        unused-slot handling (e.g. `_maybe_split_business_process_slide`)
        omit `slide` to skip this.
        """
        for i, pair in enumerate(pairs):
            if i >= len(items):
                # Clear unused pair shapes
                if isinstance(pair, tuple):
                    for shp in pair:
                        ShapeDiscovery.clear_text(shp)
                    if slide is not None:
                        self._remove_decorations_around(slide, list(pair))
                continue
            item = items[i]
            title = getattr(item, 'title', None) or (item.get('title') if isinstance(item, dict) else '')
            body = getattr(item, 'body', None) or getattr(item, 'content', None) \
                or (item.get('body') if isinstance(item, dict) else '') \
                or (item.get('content') if isinstance(item, dict) else '')
            if not isinstance(pair, tuple):
                continue
            title_shape, body_shape = pair
            if title_shape is not None and title_shape.has_text_frame:
                # Full text, never truncated with an ellipsis — a heading cut
                # off mid-word ("Referral trac…") reads as meaningless. Instead
                # grow the box to fit however many lines it wraps to, and push
                # the body shape down by the same amount so it doesn't overlap
                # (no auto-fit here — see body's comment below for why).
                title_text = title or ''
                extra_h = self._estimate_extra_wrap_height(title_shape, title_text, font_pt=10)
                if extra_h > 0:
                    title_shape.height += extra_h
                    if body_shape is not None:
                        body_shape.top += extra_h
                title_shape.text_frame.word_wrap = True
                self._set_text_preserving_format(title_shape.text_frame, title_text)
                self._apply_text_font(title_shape.text_frame)
                self._set_vertical_anchor_top(title_shape.text_frame)
            if body_shape is not None and body_shape.has_text_frame:
                # font_pt matches _apply_text_font's default (10pt) below — otherwise
                # the char-count estimate (previously assuming 8pt) under-truncates
                # relative to what actually gets rendered, and on template cells with
                # very little height to begin with (some of this template's cells are
                # barely 0.13" tall) the overflow renders as overlapping text lines.
                # No auto-fit here: this template's cells carry ~0.23" top+bottom
                # text insets against boxes as short as 0.13-0.34" tall, so
                # normAutofit's fit calculation goes deeply negative and some
                # renderers (LibreOffice) respond by shrinking the text to
                # invisibility rather than clamping to a sane minimum. Truncating
                # to the estimated capacity is enough to keep text on-shape.
                max_chars = self._estimate_max_chars(body_shape, font_pt=10)
                body_text = self._truncate_at_sentence(body or '', max_chars=max_chars)
                body_shape.text_frame.word_wrap = True
                self._set_text_preserving_format(body_shape.text_frame, body_text)
                self._apply_text_font(body_shape.text_frame)
                self._set_vertical_anchor_top(body_shape.text_frame)

    def _fill_table_role(
        self, prs: Presentation, slide, slide_index: int,
        shapes: list, data, role: ShapeRole,
    ) -> int:
        """Fill a table role. Reuses overflow logic by converting to markdown."""
        if not shapes or shapes[0] is None:
            print(f"  - Table shape not found for role '{role.name}'")
            return 0
        shape = shapes[0]
        rows = data if isinstance(data, list) else []
        if not rows:
            return 0
        md_table = self._rows_to_md_table(rows, fill_cols=role.fill_cols)
        target = _LegacyTarget(content_key=role.name, fill_cols=role.fill_cols)
        # Normalize column widths before estimating — otherwise the estimate
        # uses the template's (possibly narrower) placeholder widths while the
        # actual fill later normalizes them, causing a small mismatch between
        # the estimated and real wrap-line counts. Safe to call early: it only
        # touches column widths, not rows, so the later call in _fill_table is
        # a no-op once already normalized.
        if shape.has_table:
            self._normalize_table_col_widths(shape.table)
        # Use shape-position-aware, content-aware row limit so tables never
        # overflow into the footer (estimated from this exact data, not the
        # template's placeholder text). fill_cols mode emits row-only markdown
        # (no header line), so every parsed row is data; otherwise the first
        # parsed row is the header.
        parsed_rows = self._parse_md_table_rows(md_table)
        if role.fill_cols:
            data_rows = parsed_rows
        else:
            data_rows = parsed_rows[1:] if len(parsed_rows) > 1 else []
        max_rows = self._max_rows_for_shape(shape, rows_data=data_rows)
        chunks = self._split_table_data(md_table, max_rows)
        if len(chunks) <= 1:
            self._fill_table(shape, md_table, target)
            return 0

        print(f'  - Role "{role.name}": table split ({len(chunks)} slides, max {max_rows} rows/slide)')
        self._fill_table(shape, chunks[0], target)
        extra = 0
        current_index = slide_index
        for i, chunk in enumerate(chunks[1:], 1):
            new_slide = self._duplicate_slide(prs, current_index)
            extra += 1
            current_index += 1
            # Re-discover same shape on duplicate, then fill
            new_shape = ShapeDiscovery().discover(new_slide, role)
            if new_shape and new_shape[0] is not None:
                self._fill_table(new_shape[0], chunk, target)
                print(f'  - Filled continuation {i}/{len(chunks)-1}')
        return extra

    def _fill_image_role(self, slide, shapes: list, data, role: ShapeRole):
        """Fill an image role. data: str (image path or data URI)."""
        if not shapes or shapes[0] is None:
            print(f"  - Image shape not found for role '{role.name}'")
            return
        if not data:
            return
        if getattr(role, 'fit_to_slide', False):
            self._place_image_fit_slide(slide, data, role.name)
        else:
            target = _LegacyTarget(content_key=role.name)
            self._replace_shape_with_image(slide, shapes[0], data, target)

    def _place_image_fit_slide(self, slide, image_path: str, role_name: str):
        """Scale image to fill slide content area (below title), preserving aspect ratio.

        Uses Pillow to read native image dimensions.  Falls back to placeholder-based
        placement when the file cannot be opened.
        """
        image_path = image_path.strip()

        # Resolve relative path
        if not Path(image_path).is_absolute():
            for candidate in [Path(image_path),
                               Path('outputs') / image_path,
                               (self._output_dir / image_path) if self._output_dir else None]:
                if candidate and candidate.exists():
                    image_path = str(candidate)
                    break

        if not Path(image_path).exists():
            print(f'  - Image not found for fit-to-slide: {image_path}')
            return

        try:
            from PIL import Image as _PILImage
            with _PILImage.open(image_path) as img:
                img_w, img_h = img.size
        except Exception as e:
            print(f'  - Could not read image dimensions ({e}); skipping fit-to-slide')
            return

        # Use actual slide dimensions recorded during render_from_profile.
        # Fall back to common 10" width if not yet set.
        slide_w = self._slide_width_emu or 9144000
        # Use the dynamically computed safe-bottom (above footer) as the effective
        # lower boundary so the image never overlaps the footer band.
        safe_bottom = self._SVN_SAFE_BOTTOM_EMU
        margin_l = int(0.5 * self._EMU_PER_INCH)   # 0.5"
        margin_r = int(0.5 * self._EMU_PER_INCH)   # 0.5"
        margin_t = int(1.1 * self._EMU_PER_INCH)   # 1.1" (below title bar)

        avail_w = slide_w - margin_l - margin_r
        avail_h = max(0, safe_bottom - margin_t)

        scale = min(avail_w / img_w, avail_h / img_h)
        new_w = int(img_w * scale)
        new_h = int(img_h * scale)

        # Center within the content area
        left = margin_l + (avail_w - new_w) // 2
        top = margin_t + (avail_h - new_h) // 2

        try:
            slide.shapes.add_picture(image_path, left, top, width=new_w, height=new_h)
            print(f"  Placed '{role_name}' fit-to-slide: "
                  f"{new_w / 914400:.2f}\" × {new_h / 914400:.2f}\" "
                  f"at ({left / 914400:.2f}\", {top / 914400:.2f}\")")
        except Exception as e:
            print(f'  - Error placing fit-to-slide image: {e}')

    # ------------------------------------------------------------------
    # Helpers for role-based fill
    # ------------------------------------------------------------------

    @staticmethod
    def _coerce_text_items(data) -> list:
        """Normalize text-role data into a list of strings."""
        if data is None:
            return []
        if isinstance(data, str):
            return [data]
        if isinstance(data, list):
            return [str(item) if item is not None else '' for item in data]
        return [str(data)]

    @staticmethod
    def _rows_to_md_table(rows: list, fill_cols=None) -> str:
        """Convert list[dict | dataclass] to markdown pipe-table string.

        - When fill_cols is None: emit standard header + data rows.
        - When fill_cols is set: emit row-only format (no header/separator), and
          emit only the *last* len(fill_cols) fields per row (template provides
          fixed columns for the rest, e.g. assumption labels).
        """
        if not rows:
            return ''
        # Normalize rows to list[dict]
        import dataclasses as _dc
        rows = [r if isinstance(r, dict) else _dc.asdict(r) if _dc.is_dataclass(r) else dict(vars(r))
                for r in rows]
        headers = list(rows[0].keys())

        def esc(v):
            return str(v).replace('|', '&#124;').replace('\n', ' ')

        if fill_cols:
            n = len(fill_cols)
            target_keys = headers[-n:] if n <= len(headers) else headers
            lines = []
            for r in rows:
                cells = [esc(r.get(k, '')) for k in target_keys]
                lines.append('| ' + ' | '.join(cells) + ' |')
            return '\n'.join(lines)

        lines = [
            '| ' + ' | '.join(headers) + ' |',
            '|' + '|'.join(['---'] * len(headers)) + '|',
        ]
        for r in rows:
            lines.append('| ' + ' | '.join(esc(r.get(h, '')) for h in headers) + ' |')
        return '\n'.join(lines)


    # ==================================================================
    # Extra-slide rendering — AI-generated, generate-slide-style layouts
    # ==================================================================

    def _render_extra_slides(
        self, prs: Presentation, extras: List[dict], last_positions: Dict[int, int],
        blank_canvas_layout=None,
    ) -> None:
        """Insert AI-generated slides rendered with generate-slide-style layouts.

        Each entry of *extras* is a dict with keys:
          - layout:          str  (numbered_points|card_grid|comparison_2|... ; default bullets)
          - title:           str  (slide title)
          - <layout content> (e.g. points / items / columns / steps / message)
          - anchor_section:  str | None (profile section name; resolves to anchor slide)
          - anchor_slide:    int | None (explicit slide number; overrides anchor_section)

        New slides are created from a CONTENT slide's layout (so they inherit the
        deck's white background + master footer/page-number/logo) and inserted
        right after the LAST variant of the anchor slide (accounting for overflow
        inserts). Multiple extras sharing an anchor keep their given order.
        """
        from lib.templates.svn import SECTION_TO_SLIDES

        w_emu, h_emu = prs.slide_width, prs.slide_height

        # Prefer a chrome reference slide — a configured CONTENT slide that
        # has slide-level chrome shapes (logo / footer in the bottom strip).
        # Cloning that slide makes the chrome real slide-level shapes (editable)
        # rather than master-only chrome (which appears "captured" / locked).
        chrome_ref_idx = self._chrome_canvas_idx(prs, last_positions)
        canvas_layout = None
        if chrome_ref_idx is None:
            # Fallback: master-chrome only (legacy behavior).
            canvas_layout = self._content_canvas_layout(prs, last_positions, blank_canvas_layout)
            if canvas_layout is None:
                print('No content-slide layout available for extra slides')
                return

        # Group by anchor (preserve given order within a group)
        groups: Dict[int, List[dict]] = {}
        for e in extras:
            anchor = self._resolve_extra_anchor(e, last_positions, SECTION_TO_SLIDES, prs)
            groups.setdefault(anchor, []).append(e)

        # Insertions shift later positions; process in ascending anchor order
        # and track cumulative offset.
        offset = 0
        for anchor in sorted(groups.keys()):
            for i, extra in enumerate(groups[anchor]):
                target = anchor + offset + 1  # 0-indexed insertion point
                if chrome_ref_idx is not None:
                    # Source index for chrome_ref shifts right with every prior
                    # insert that landed at or before it.
                    src_idx = chrome_ref_idx
                    new_slide = self._duplicate_slide_at(prs, src_idx, target)
                    self._strip_content_keep_chrome(new_slide, h_emu)
                    self._fill_extra_slide(new_slide, extra, w_emu, h_emu, clear=False)
                    if src_idx >= target:
                        chrome_ref_idx += 1  # our insert shifted the ref slide
                else:
                    new_slide = self._add_slide_at(prs, canvas_layout, target)
                    self._fill_extra_slide(new_slide, extra, w_emu, h_emu)
                offset += 1
                print(f'  - Inserted extra slide "{extra.get("title", "")[:40]}"'
                      f' [{extra.get("layout", "bullets")}] at index {target}')

    def _content_canvas_layout(self, prs: Presentation, last_positions: Dict[int, int],
                               blank_canvas_layout=None):
        """Pick a CLEAN content-slide layout (light bg + standard small footer logo).

        Some configured slides (e.g. slide 4 in the SVN deck) use a unique master
        with a large decorative side panel and an oversized logo. Using that
        layout for AI-generated extra slides makes them visually inconsistent
        with the rest of the deck. We skip those and prefer the layout of a
        configured content slide whose master only contains small chrome
        (footer text, page number, small logo).
        """
        if not last_positions:
            return blank_canvas_layout or (prs.slides[0].slide_layout if len(prs.slides) else None)

        # Iterate configured slides in template-order; pick the first clean one.
        slide_w = prs.slide_width or 0
        slide_h = prs.slide_height or 0
        slide_area = slide_w * slide_h
        for sn in sorted(last_positions):
            idx = last_positions[sn]
            if not (0 <= idx < len(prs.slides)):
                continue
            layout = prs.slides[idx].slide_layout
            if self._is_clean_master(layout.slide_master, slide_area):
                return layout

        # None of the surviving configured slides has a clean master (e.g. only
        # a two-tone divider-style slide like project_background survived
        # pruning). Reusing its layout would bake that decorative background
        # into every extra slide, so prefer the blank end-card layout captured
        # before pruning instead.
        if blank_canvas_layout is not None:
            return blank_canvas_layout

        # Last resort: earliest configured slide's layout (legacy behavior)
        idx = last_positions[min(last_positions)]
        if 0 <= idx < len(prs.slides):
            return prs.slides[idx].slide_layout
        return prs.slides[0].slide_layout if len(prs.slides) else None

    @staticmethod
    def _is_clean_master(master, slide_area: int) -> bool:
        """A master is 'clean' if it has no large solid-filled decorative shapes.

        Large = >= 25% of the slide area. This filters out section-divider style
        masters (e.g. SVN slide 4 with a pink right-half panel) while keeping
        standard content-slide masters that only carry footer / logo chrome.
        """
        if slide_area <= 0:
            return True
        threshold = slide_area * 0.25
        try:
            for sh in master.shapes:
                if not (sh.width and sh.height):
                    continue
                if sh.width * sh.height < threshold:
                    continue
                # Large shape — clean only if it's not a solid colored fill
                try:
                    fill = sh.fill
                    if fill.type == 1:  # MSO_FILL.SOLID
                        return False
                except Exception:
                    pass
        except Exception:
            return True
        return True

    def _add_slide_at(self, prs: Presentation, layout, target_index: int):
        """Add a fresh slide from *layout* and move it to target_index (0-based)."""
        new_slide = prs.slides.add_slide(layout)
        self._ensure_unique_slide_partname(prs, new_slide)
        pml_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        try:
            sldIdLst = prs.part._element.find(f'{{{pml_ns}}}sldIdLst')
            if sldIdLst is not None:
                new_sldId = list(sldIdLst)[-1]  # just appended → at end
                sldIdLst.remove(new_sldId)
                remaining = list(sldIdLst)
                if target_index >= len(remaining):
                    sldIdLst.append(new_sldId)
                else:
                    remaining[target_index].addprevious(new_sldId)
        except Exception as e:
            print(f'  - Error positioning extra slide: {e}')
        return prs.slides[target_index]

    def _resolve_extra_anchor(
        self, extra: dict, last_positions: Dict[int, int],
        section_map: Dict[str, list], prs: Presentation,
    ) -> int:
        """Return 0-indexed slide position to insert AFTER.

        ``last_positions`` reflects POST-prune positions — a requested anchor
        (explicit slide number or section) may no longer exist if it was
        pruned for having no content. In that case fall back to appending
        after the last surviving slide rather than using a stale/incorrect
        index (there is no valid "pre-prune index" once slides are deleted).
        """
        fallback = max(last_positions.values()) if last_positions else len(prs.slides) - 2

        explicit = extra.get('anchor_slide')
        if isinstance(explicit, int) and explicit > 0:
            return last_positions.get(explicit, fallback)
        section = extra.get('anchor_section')
        if section and section in section_map:
            slide_num = section_map[section][-1]
            return last_positions.get(slide_num, fallback)
        # Default: append before the extra-slide template (i.e. after last filled slide)
        if last_positions:
            return max(last_positions.values())
        return len(prs.slides) - 2  # before the template slide itself

    def _fill_extra_slide(self, slide, extra: dict, w_emu: int, h_emu: int,
                          clear: bool = True) -> None:
        """Render a generate-slide-style layout onto the cloned canvas.

        When *clear* is True the slide is cleared first (use for slides created
        from a layout). When False, existing shapes are kept — used after the
        slide was duplicated from an anchor and its content shapes stripped, so
        only the chrome (logo / footer) remains and the layout is drawn on top.

        See lib/extra_slide_layouts.py for supported layouts.
        """
        from lib.extra_slide_layouts import render_extra_layout
        render_extra_layout(slide, extra, w_emu, h_emu, clear=clear)

    # ------------------------------------------------------------------
    # Chrome-aware slide cloning helpers (used for extra AI-generated slides)
    # ------------------------------------------------------------------

    def _chrome_canvas_idx(self, prs: Presentation,
                           last_positions: Dict[int, int]) -> Optional[int]:
        """Find a configured content slide that has slide-level chrome shapes.

        Chrome shapes are small elements parked in the bottom strip of the
        slide (logo, footer text). Cloning such a slide preserves the chrome as
        real slide-level shapes — editable in PowerPoint, rather than appearing
        as locked master-only chrome.

        Returns the slide index, or None when no suitable slide is found.
        """
        if not last_positions:
            return None
        slide_w = prs.slide_width or 0
        slide_h = prs.slide_height or 0
        if slide_h <= 0:
            return None
        slide_area = slide_w * slide_h
        threshold_top = int(slide_h * 0.88)  # shapes anchored in bottom ~12%
        # Iterate configured slides in template order, prefer clean masters.
        for sn in sorted(last_positions):
            idx = last_positions[sn]
            if not (0 <= idx < len(prs.slides)):
                continue
            slide = prs.slides[idx]
            if not self._is_clean_master(slide.slide_layout.slide_master, slide_area):
                continue
            for sh in slide.shapes:
                if sh.top is not None and sh.top >= threshold_top:
                    return idx
        return None

    def _duplicate_slide_at(self, prs: Presentation, source_idx: int,
                            target_idx: int):
        """Duplicate slide at *source_idx* and insert the copy at *target_idx*.

        Differs from :meth:`_duplicate_slide` (which always inserts immediately
        after the source) in that it places the copy at an arbitrary position.
        Also re-binds any relationship references (``r:embed`` / ``r:link``
        / ``r:id``) inside the cloned shapes so embedded images survive the
        cross-slide copy — otherwise the new slide's blipFills would resolve
        their rIds against the wrong relationship table.
        """
        source_slide = prs.slides[source_idx]

        new_slide = prs.slides.add_slide(source_slide.slide_layout)
        self._ensure_unique_slide_partname(prs, new_slide)
        src_tree = source_slide.shapes._spTree
        new_tree = new_slide.shapes._spTree
        for child in list(new_tree):
            new_tree.remove(child)
        for child in src_tree:
            new_tree.append(copy.deepcopy(child))

        self._remap_shape_rels(source_slide.part, new_slide.part, new_tree)

        pml_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        try:
            sldIdLst = prs.part._element.find(f'{{{pml_ns}}}sldIdLst')
            if sldIdLst is not None:
                new_sldId = list(sldIdLst)[-1]  # just appended at end
                sldIdLst.remove(new_sldId)
                remaining = list(sldIdLst)
                if target_idx >= len(remaining):
                    sldIdLst.append(new_sldId)
                else:
                    remaining[target_idx].addprevious(new_sldId)
        except Exception as e:
            print(f'  - Error positioning extra slide: {e}')
        return prs.slides[target_idx]

    @staticmethod
    def _remap_shape_rels(source_part, target_part, new_spTree) -> None:
        """Re-bind ``r:embed`` / ``r:link`` / ``r:id`` references in *new_spTree*
        to relationships on *target_part*.

        Cross-slide deep-copies preserve the textual rId values (e.g. ``rId1``)
        but those IDs are scoped to the originating slide's relationship table.
        Without remapping, an image blipFill copied from another slide ends up
        pointing at whatever the destination slide's ``rId1`` happens to be
        (typically the slideLayout), causing the image to vanish.
        """
        r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        src_rels = source_part.rels
        rid_map: Dict[str, str] = {}
        for el in new_spTree.iter():
            for attr_name in list(el.attrib):
                if not attr_name.startswith(f'{{{r_ns}}}'):
                    continue
                local = attr_name.split('}', 1)[1]
                if local not in ('embed', 'link', 'id'):
                    continue
                old_rid = el.attrib[attr_name]
                if not old_rid:
                    continue
                if old_rid in rid_map:
                    el.attrib[attr_name] = rid_map[old_rid]
                    continue
                try:
                    rel = src_rels[old_rid]
                except KeyError:
                    continue
                try:
                    if rel.is_external:
                        new_rid = target_part.relate_to(
                            rel.target_ref, rel.reltype, is_external=True)
                    else:
                        new_rid = target_part.relate_to(
                            rel.target_part, rel.reltype)
                except Exception as exc:
                    print(f'  - Warn: could not rebind {old_rid} ({rel.reltype}): {exc}')
                    continue
                rid_map[old_rid] = new_rid
                el.attrib[attr_name] = new_rid

    @staticmethod
    def _strip_content_keep_chrome(slide, slide_h_emu: int) -> None:
        """Remove every shape whose top edge sits above the bottom chrome strip.

        Anything anchored in the bottom ~12% of the slide is treated as chrome
        (logo / footer text) and retained so the layout renderer can draw the
        new content on top while the chrome remains intact and editable.
        """
        if slide_h_emu <= 0:
            return
        threshold = int(slide_h_emu * 0.88)
        for sh in list(slide.shapes):
            if sh.top is None or sh.top < threshold:
                sh._element.getparent().remove(sh._element)


class _LegacyTarget:
    """Lightweight stand-in for legacy ShapeTarget used by _fill_table / _replace_shape_with_image."""

    def __init__(self, content_key: str, fill_cols=None):
        self.content_key = content_key
        self.fill_cols = fill_cols
        self.shape_index = None
        self.shape_name = None
